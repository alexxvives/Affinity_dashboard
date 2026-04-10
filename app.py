# app.py — Segment Effect Explorer (v2)
# Run: streamlit run app.py

import ast
import io
import os
from concurrent.futures import ThreadPoolExecutor
from pathlib import Path
from typing import Dict, List, Optional, Tuple

_APP_DIR = Path(__file__).parent

import numpy as np
import pandas as pd
import plotly.express as px
import streamlit as st
import streamlit.components.v1 as components

SEGMENT_LABELS: Dict[str, str] = {}
SEGMENT_DESCRIPTIONS: Dict[str, str] = {}
PRODUCT_COLS: List[str] = [
    "Business Line Of Credit", "Cdira", "Checking", "Commercial Loan",
    "Credit Card", "Escrow", "Standalone Savings", "Home Equity",
    "Investments", "Loan", "Loan - Personal", "Loc - Personal",
    "Money Market", "Mortgage", "Odloc", "Other", "Savings",
]

# Load real segment data from segment_descriptions.csv (overrides generate_dummy_data values)
try:
    _sdesc_df = pd.read_csv(_APP_DIR / "segment_descriptions.csv", dtype=str)
    if {"nsegment", "label", "description"}.issubset(_sdesc_df.columns):
        SEGMENT_LABELS = dict(zip(_sdesc_df["nsegment"], _sdesc_df["label"].fillna("")))
        SEGMENT_DESCRIPTIONS = dict(zip(_sdesc_df["nsegment"], _sdesc_df["description"].fillna("")))
        _SEGMENT_GROUPS: Dict[str, List[str]] = {}
        for _, _srow in _sdesc_df.iterrows():
            _SEGMENT_GROUPS.setdefault(_srow["label"], []).append(_srow["nsegment"])
    else:
        _SEGMENT_GROUPS: Dict[str, List[str]] = {}
except FileNotFoundError:
    _SEGMENT_GROUPS: Dict[str, List[str]] = {}

# ── Config ────────────────────────────────────────────────────────────────────
COMM_ORDER = ["day1", "day5", "day7", "day31", "day61DD", "day61NDD", "day90", "day120"]
REQUIRED_COLS = [
    "communication", "alpha_key", "contact_flag",
    "start_date", "end_date",
    "start_balance", "end_balance",
    "start_accounts", "end_accounts",
    "nsegments",
]

Z95 = 1.96   # z-score for 95% CI


# ── Helpers ───────────────────────────────────────────────────────────────────
def _try_parse_listlike(x):
    if x is None or (isinstance(x, float) and np.isnan(x)):
        return []
    if isinstance(x, list):
        return [str(v) for v in x if v is not None and str(v).strip() != ""]
    if isinstance(x, (tuple, set)):
        return [str(v) for v in list(x) if v is not None and str(v).strip() != ""]
    if isinstance(x, str):
        s = x.strip()
        if s == "":
            return []
        try:
            val = ast.literal_eval(s)
            if isinstance(val, list):
                return [str(v) for v in val if v is not None and str(v).strip() != ""]
        except Exception:
            pass
        if "," in s:
            parts = [p.strip().strip(chr(39)).strip(chr(34)) for p in s.split(",")]
            return [p for p in parts if p != ""]
        # dash-delimited format: "-000302-3290392-3203203-" → ["000302", "3290392", "3203203"]
        if s.startswith("-") and s.endswith("-") and len(s) > 1:
            parts = [p.strip() for p in s.strip("-").split("-")]
            return [p for p in parts if p != ""]
        return [s.strip(chr(39)).strip(chr(34))]
    return [str(x)]


def _df_hash(df: pd.DataFrame) -> int:
    """Fast approximate hash for large DataFrames — samples ~1 000 evenly-spaced rows."""
    step = max(1, len(df) // 1000)
    return int(pd.util.hash_pandas_object(df.iloc[::step]).sum()) ^ hash(df.shape)


@st.cache_data(show_spinner=False)
def _load_csv(b: bytes) -> pd.DataFrame:
    df = pd.read_csv(io.BytesIO(b))
    return df.rename(columns={"Communication": "communication", "Contact_flag": "contact_flag"})


@st.cache_data(show_spinner=False)
def _load_audience_profile(b: bytes) -> pd.DataFrame:
    """Load audience_profile.csv from bytes, pre-computing derived columns if not already present."""
    aud = pd.read_csv(io.BytesIO(b))
    today = pd.Timestamp.today().normalize()
    if "age" not in aud.columns or not pd.api.types.is_numeric_dtype(aud["age"]):
        aud["age"] = (today - pd.to_datetime(aud["date_of_birth"], errors="coerce")).dt.days / 365.25
    if "tenure_years" not in aud.columns or not pd.api.types.is_numeric_dtype(aud["tenure_years"]):
        aud["tenure_years"] = (today - pd.to_datetime(aud["date_first_relation"], errors="coerce")).dt.days / 365.25
    if "sow" not in aud.columns:
        aud["sow"] = (aud["amount_deposit_spot_balance"] / aud["total_deposits_ixi"]).clip(0, 1)
    if "n_products" not in aud.columns:
        present_prod_cols = [c for c in PRODUCT_COLS if c in aud.columns]
        aud["n_products"] = aud[present_prod_cols].sum(axis=1).astype(int)
    return aud


# ── CC BT campaign helpers ────────────────────────────────────────────────────

CAMPAIGNS: Dict[str, Dict] = {
    "SELECTCHK": {"file": "SELECTCHK_campaign.csv", "label": "SELECT CHK"},
    "CC_BT":     {"file": "CC_BT_campaign.csv",     "label": "CC Balance Transfer"},
}


@st.cache_data(show_spinner=False)
def _load_cc_bt(b: bytes) -> pd.DataFrame:
    df = pd.read_csv(io.BytesIO(b))
    df.columns = df.columns.str.lower().str.strip()
    return df


@st.cache_data(show_spinner=False, hash_funcs={pd.DataFrame: _df_hash})
def preprocess_cc_bt(df_raw: pd.DataFrame,
                     date_min: Optional[str] = None,
                     date_max: Optional[str] = None) -> pd.DataFrame:
    df = df_raw.copy()
    if "bt_date" in df.columns and (date_min or date_max):
        df["bt_date"] = pd.to_datetime(df["bt_date"], errors="coerce")
        if date_min:
            df = df[df["bt_date"] >= pd.Timestamp(date_min)]
        if date_max:
            df = df[df["bt_date"] <= pd.Timestamp(date_max)]
    df["nsegments"] = df["nsegments"].apply(_try_parse_listlike)
    df["nsegments"] = df["nsegments"].apply(lambda xs: xs if xs else ["__NO_SEGMENT__"])
    df = df.explode("nsegments").rename(columns={"nsegments": "nsegment"})
    df["nsegment"] = df["nsegment"].astype(str)
    return df


@st.cache_data(show_spinner=False, hash_funcs={pd.DataFrame: _df_hash})
def agg_cc_bt(df_bt: pd.DataFrame, min_n: int = 30) -> pd.DataFrame:
    """Per-segment BT conversion rate lift and BT amount lift."""
    treated = df_bt[df_bt["control_flag"] == 0]
    control = df_bt[df_bt["control_flag"] == 1]

    g_t = treated.groupby("nsegment")
    g_c = control.groupby("nsegment")

    n_t = g_t["alpha_key"].count()
    n_c = g_c["alpha_key"].count().reindex(n_t.index, fill_value=0)

    mean_flag_t = g_t["bt_flag"].mean()
    mean_flag_c = g_c["bt_flag"].mean().reindex(mean_flag_t.index)
    mean_amt_t  = g_t["bt_amount"].mean()
    mean_amt_c  = g_c["bt_amount"].mean().reindex(mean_amt_t.index)

    std_ft = g_t["bt_flag"].std(ddof=0)
    std_fc = g_c["bt_flag"].std(ddof=0).reindex(mean_flag_t.index)
    std_at = g_t["bt_amount"].std(ddof=0)
    std_ac = g_c["bt_amount"].std(ddof=0).reindex(mean_flag_t.index)

    se_ft = (std_ft / np.sqrt(n_t)).where(n_t >= 2)
    se_fc = (std_fc / np.sqrt(n_c)).where(n_c >= 2)
    se_at = (std_at / np.sqrt(n_t)).where(n_t >= 2)
    se_ac = (std_ac / np.sqrt(n_c)).where(n_c >= 2)

    out = pd.DataFrame({
        "nsegment":       mean_flag_t.index,
        "n_treated":      n_t.values,
        "n_control":      n_c.values,
        "conv_treated":   mean_flag_t.values,
        "conv_control":   mean_flag_c.values,
        "conv_lift":      (mean_flag_t - mean_flag_c).values,
        "conv_lift_ci":   (Z95 * np.sqrt(se_ft**2 + se_fc**2)).values,
        "amt_treated":    mean_amt_t.values,
        "amt_control":    mean_amt_c.values,
        "amt_lift":       (mean_amt_t - mean_amt_c).values,
        "amt_lift_ci":    (Z95 * np.sqrt(se_at**2 + se_ac**2)).values,
    })
    out = out[out["n_treated"].fillna(0) >= min_n]
    return out.set_index("nsegment")


@st.cache_data(show_spinner=False)
def _cc_bt_seg_ids(df_raw: pd.DataFrame) -> List[str]:
    segs: set = set()
    for v in df_raw["nsegments"]:
        for s in _try_parse_listlike(v):
            if s and s != "__NO_SEGMENT__":
                segs.add(s)
    return sorted(segs)


@st.cache_data(show_spinner=False, hash_funcs={pd.DataFrame: _df_hash})
def preprocess(df_raw: pd.DataFrame, date_min: Optional[str], date_max: Optional[str], bal_clip_pct: float = 0.0) -> pd.DataFrame:
    df = df_raw.copy()
    df.columns = df.columns.str.lower().str.strip()  # normalise: Communication→communication, Contact_flag→contact_flag
    df["start_date"] = pd.to_datetime(df["start_date"], errors="coerce")
    df["end_date"]   = pd.to_datetime(df["end_date"],   errors="coerce")
    df["communication"] = df["communication"].astype(str).str.strip()
    df["contact_flag"]  = pd.to_numeric(df["contact_flag"], errors="coerce").fillna(0).astype(int)
    df["control_flag"] = (df["contact_flag"] == 0).astype(int)
    for col in ["start_balance", "end_balance", "start_accounts", "end_accounts"]:
        df[col] = pd.to_numeric(df[col], errors="coerce")

    # Date range filter
    if date_min:
        df = df[df["start_date"] >= pd.Timestamp(date_min)]
    if date_max:
        df = df[df["start_date"] <= pd.Timestamp(date_max)]

    # Normalise dash-delimited strings like "-000302-3290392-" into a list before further parsing
    df["nsegments"] = df["nsegments"].apply(_try_parse_listlike)
    df["nsegments"] = df["nsegments"].apply(lambda xs: xs if xs else ["__NO_SEGMENT__"])
    df = df.explode("nsegments").rename(columns={"nsegments": "nsegment"})
    df["nsegment"] = df["nsegment"].astype(str)

    eps = 1e-6
    df["balance_abs_change"] = df["end_balance"] - df["start_balance"]
    denom_b = df["start_balance"].abs().replace(0, np.nan)
    df["balance_pct_change"] = df["balance_abs_change"] / denom_b
    df.loc[df["start_balance"].abs() < eps, "balance_pct_change"] = np.nan

    df["accounts_abs_change"] = df["end_accounts"] - df["start_accounts"]
    denom_a = df["start_accounts"].replace(0, np.nan)
    df["accounts_pct_change"] = df["accounts_abs_change"] / denom_a
    df.loc[df["start_accounts"].fillna(0) == 0, "accounts_pct_change"] = np.nan

    # Outlier clipping on balance % change
    if bal_clip_pct > 0:
        lo = df["balance_pct_change"].quantile(bal_clip_pct / 100.0)
        hi = df["balance_pct_change"].quantile(1.0 - bal_clip_pct / 100.0)
        df["balance_pct_change"] = df["balance_pct_change"].clip(lo, hi)

    return df


def _rank(c: str) -> int:
    try:
        return COMM_ORDER.index(c)
    except ValueError:
        return 999


@st.cache_data(show_spinner=False)
def _all_segment_ids(df_raw: pd.DataFrame) -> List[str]:
    segs: set = set()
    for v in df_raw["nsegments"]:
        for s in _try_parse_listlike(v):
            if s and s != "__NO_SEGMENT__":
                segs.add(s)
    return sorted(segs)


@st.fragment
def _seg_lookup_widget(seg_ids: List[str], labels: Dict[str, str], descriptions: Dict[str, str]) -> None:
    """Searchable segment lookup — Streamlit native search handles filtering as the user types."""
    sel = st.selectbox(
        "Segment lookup",
        options=seg_ids,
        index=None,
        placeholder="Type segment ID or name…",
        key="seg_lookup_sel",
        format_func=lambda x: f"{x}  —  {labels.get(x, '')}",
    )
    if sel:
        lbl  = labels.get(sel, "")
        desc = descriptions.get(sel, "")
        if lbl:
            st.caption(f"**{lbl}** — {desc}" if desc else f"**{lbl}**")
        elif desc:
            st.caption(desc)
    st.divider()


# ── Aggregation ───────────────────────────────────────────────────────────────
@st.cache_data(show_spinner=False, hash_funcs={pd.DataFrame: _df_hash})
def agg_cols_data(
    df: pd.DataFrame,
    selected: Tuple[str, ...],
    all_mode: bool,
    bal_baseline_min: Optional[float],
) -> pd.DataFrame:
    EMPTY = pd.DataFrame(columns=["nsegment", "agg_bal", "agg_bal_ci", "agg_acct", "agg_acct_ci", "agg_lift_bal", "agg_lift_bal_ci", "agg_lift_acct", "agg_lift_acct_ci", "agg_n"])
    if not selected:
        return EMPTY
    treated = df[(df["contact_flag"] == 1) & (df["communication"].isin(list(selected)))].copy()
    control = df[(df["control_flag"] == 1) & (df["communication"].isin(list(selected)))].copy()
    if treated.empty:
        return EMPTY
    treated["_r"] = treated["communication"].map(_rank)
    if all_mode:
        cnt = (treated.groupby(["alpha_key", "nsegment"])["communication"]
               .nunique().reset_index(name="_c"))
        ok = cnt[cnt["_c"] >= len(set(selected))][["alpha_key", "nsegment"]]
        treated = treated.merge(ok, on=["alpha_key", "nsegment"], how="inner")
        if treated.empty:
            return EMPTY
    treated = (treated.sort_values("_r", ascending=False)
                      .drop_duplicates(subset=["alpha_key", "nsegment"], keep="first"))
    # Low-balance filter: exclude segments with mean start_balance below minimum
    if bal_baseline_min is not None:
        seg_base = treated.groupby("nsegment")["start_balance"].mean()
        ok_segs = seg_base[seg_base >= bal_baseline_min].index
        treated = treated[treated["nsegment"].isin(ok_segs)]
    if treated.empty:
        return EMPTY

    # ── Vectorised per-segment stats ─────────────────────────────────────────
    g_t      = treated.groupby("nsegment", dropna=False)
    cnt_b_t  = g_t["balance_pct_change"].count()
    cnt_a_t  = g_t["accounts_pct_change"].count()
    mean_b_t = g_t["balance_pct_change"].mean()
    mean_a_t = g_t["accounts_pct_change"].mean()
    std_b_t  = g_t["balance_pct_change"].std(ddof=0)
    std_a_t  = g_t["accounts_pct_change"].std(ddof=0)
    se_b_t   = (std_b_t / np.sqrt(cnt_b_t)).where(cnt_b_t >= 2)
    se_a_t   = (std_a_t / np.sqrt(cnt_a_t)).where(cnt_a_t >= 2)
    n_u      = g_t["alpha_key"].nunique()

    if not control.empty:
        g_c      = control.groupby("nsegment", dropna=False)
        mean_b_c = g_c["balance_pct_change"].mean().reindex(mean_b_t.index)
        mean_a_c = g_c["accounts_pct_change"].mean().reindex(mean_a_t.index)
        cnt_b_c  = g_c["balance_pct_change"].count().reindex(mean_b_t.index)
        cnt_a_c  = g_c["accounts_pct_change"].count().reindex(mean_a_t.index)
        std_b_c  = g_c["balance_pct_change"].std(ddof=0).reindex(mean_b_t.index)
        std_a_c  = g_c["accounts_pct_change"].std(ddof=0).reindex(mean_a_t.index)
        se_b_c   = (std_b_c / np.sqrt(cnt_b_c)).where(cnt_b_c >= 2)
        se_a_c   = (std_a_c / np.sqrt(cnt_a_c)).where(cnt_a_c >= 2)
        lift_b    = mean_b_t - mean_b_c
        lift_a    = mean_a_t - mean_a_c
        lift_b_ci = Z95 * np.sqrt(se_b_t**2 + se_b_c**2)
        lift_a_ci = Z95 * np.sqrt(se_a_t**2 + se_a_c**2)
    else:
        lift_b = lift_a = pd.Series(np.nan, index=mean_b_t.index)
        lift_b_ci = lift_a_ci = pd.Series(np.nan, index=mean_b_t.index)

    out = pd.DataFrame({
        "nsegment":         mean_b_t.index,
        "agg_bal":          mean_b_t.values,
        "agg_bal_ci":       (Z95 * se_b_t).values,
        "agg_acct":         mean_a_t.values,
        "agg_acct_ci":      (Z95 * se_a_t).values,
        "agg_lift_bal":     lift_b.values,
        "agg_lift_bal_ci":  lift_b_ci.values,
        "agg_lift_acct":    lift_a.values,
        "agg_lift_acct_ci": lift_a_ci.values,
        "agg_n":            n_u.values,
    })
    return out if not out.empty else EMPTY


@st.cache_data(show_spinner=False, hash_funcs={pd.DataFrame: _df_hash})
def comm_data(df: pd.DataFrame, comm: str, combo_treated_keys: Optional[frozenset] = None) -> pd.DataFrame:
    EMPTY = pd.DataFrame(columns=["nsegment", f"{comm}_bal", f"{comm}_bal_ci", f"{comm}_acct", f"{comm}_acct_ci", f"{comm}_n", f"{comm}_lift_bal", f"{comm}_lift_bal_ci", f"{comm}_lift_acct", f"{comm}_lift_acct_ci"])
    if combo_treated_keys is not None:
        # Combo mode: treatment = customers who received ALL checked comms; control = those who did NOT receive this comm
        treated = df[(df["alpha_key"].isin(combo_treated_keys)) & (df["communication"] == comm)]
        control = df[(df["contact_flag"] == 0) & (df["communication"] == comm)]
    else:
        treated = df[(df["contact_flag"] == 1) & (df["communication"] == comm)]
        control = df[(df["control_flag"] == 1) & (df["communication"] == comm)]
    if treated.empty:
        return EMPTY

    def _vstat(sub):
        """Vectorised per-segment mean, biased-SE (weights always 1.0 since recency_decay=0)."""
        g      = sub.groupby("nsegment", dropna=False)
        cnt_b  = g["balance_pct_change"].count()
        cnt_a  = g["accounts_pct_change"].count()
        mean_b = g["balance_pct_change"].mean()
        mean_a = g["accounts_pct_change"].mean()
        std_b  = g["balance_pct_change"].std(ddof=0)
        std_a  = g["accounts_pct_change"].std(ddof=0)
        # SE = biased-std / sqrt(n); NaN when fewer than 2 valid observations
        se_b   = (std_b / np.sqrt(cnt_b)).where(cnt_b >= 2)
        se_a   = (std_a / np.sqrt(cnt_a)).where(cnt_a >= 2)
        return mean_b, mean_a, se_b, se_a

    t_mb, t_ma, t_se_b, t_se_a = _vstat(treated)
    n_u = treated.groupby("nsegment", dropna=False)["alpha_key"].nunique()

    if not control.empty:
        c_mb, c_ma, c_se_b, c_se_a = _vstat(control)
        lift_b    = t_mb - c_mb.reindex(t_mb.index)
        lift_a    = t_ma - c_ma.reindex(t_ma.index)
        lift_b_ci = Z95 * np.sqrt(t_se_b**2 + c_se_b.reindex(t_mb.index)**2)
        lift_a_ci = Z95 * np.sqrt(t_se_a**2 + c_se_a.reindex(t_ma.index)**2)
    else:
        lift_b = lift_a = pd.Series(np.nan, index=t_mb.index)
        lift_b_ci = lift_a_ci = pd.Series(np.nan, index=t_mb.index)

    result = pd.DataFrame({
        "nsegment":             t_mb.index,
        f"{comm}_bal":          t_mb.values,
        f"{comm}_bal_ci":       (Z95 * t_se_b).values,
        f"{comm}_acct":         t_ma.values,
        f"{comm}_acct_ci":      (Z95 * t_se_a).values,
        f"{comm}_n":            n_u.reindex(t_mb.index).values,
        f"{comm}_lift_bal":     lift_b.values,
        f"{comm}_lift_bal_ci":  lift_b_ci.values,
        f"{comm}_lift_acct":    lift_a.values,
        f"{comm}_lift_acct_ci": lift_a_ci.values,
    })
    return result if not result.empty else EMPTY


# Minimum N for a lift to be statistically meaningful.
# Below this, the standard error is too large for the lift to be trusted.
# Rule of thumb: N=30 gives ~18% standard error on a typical proportion;
# N=50 gives ~14%. We use 30 as floor but expose it as a constant.
_MIN_N_DEFAULT = 30



@st.cache_data(show_spinner=False, hash_funcs={pd.DataFrame: _df_hash})
def build_table(
    df: pd.DataFrame,
    ordered_comms: List[str],
    all_mode: bool,
    min_n: int,
    bal_baseline_min: Optional[float],
    combo_treated_keys: Optional[frozenset] = None,
    _ver: int = 2,  # bump to invalidate stale cache
) -> pd.DataFrame:
    base = pd.DataFrame({"nsegment": sorted(df["nsegment"].unique())})
    if combo_treated_keys is None:
        # Standard mode: include aggregate column
        a = agg_cols_data(df, tuple(ordered_comms), all_mode, bal_baseline_min)
        tbl = base.merge(a, on="nsegment", how="left")
        if "agg_n" in tbl.columns:
            mask = tbl["agg_n"].fillna(0) < min_n
            for col in ["agg_bal", "agg_bal_ci", "agg_acct", "agg_acct_ci",
                        "agg_lift_bal", "agg_lift_bal_ci", "agg_lift_acct", "agg_lift_acct_ci"]:
                if col in tbl.columns:
                    tbl.loc[mask, col] = np.nan
    else:
        tbl = base  # combo mode: no aggregate column, per-comm columns only
    # Fetch per-communication stats in parallel
    # Pass frozenset (hashable) so @st.cache_data can cache comm_data results properly
    with ThreadPoolExecutor() as pool:
        comm_results = list(pool.map(lambda c: comm_data(df, c, combo_treated_keys), ordered_comms))
    for comm, ca in zip(ordered_comms, comm_results):
        tbl = tbl.merge(ca, on="nsegment", how="left")
        nc = f"{comm}_n"
        if nc in tbl.columns:
            mask = tbl[nc].fillna(0) < min_n
            for col in [f"{comm}_bal", f"{comm}_bal_ci", f"{comm}_acct", f"{comm}_acct_ci",
                        f"{comm}_lift_bal", f"{comm}_lift_bal_ci",
                        f"{comm}_lift_acct", f"{comm}_lift_acct_ci"]:
                if col in tbl.columns:
                    tbl.loc[mask, col] = np.nan
    metric_cols = [c for c in tbl.columns if c != "nsegment" and not c.endswith("_n")]
    tbl = tbl.dropna(subset=metric_cols, how="all").set_index("nsegment")
    return tbl


def _rdylgn(val: float, lo: float = -0.10, hi: float = 0.10) -> str:
    if pd.isna(val):
        return ""
    t = max(0.0, min(1.0, (val - lo) / (hi - lo)))
    if t < 0.5:
        s = t * 2
        r, g, b = 220, int(220 * s), 50
    else:
        s = (t - 0.5) * 2
        r, g, b = int(220 * (1 - s)), 200, 50
    return f"background-color: rgba({r},{g},{b},0.55); color: #111"


def style_tbl(
    tbl: pd.DataFrame,
    ordered_comms: List[str],
    seg_labels: Optional[Dict[str, str]] = None,
    seg_desc: Optional[Dict[str, str]] = None,
    show_n_cols: bool = False,
    show_lift: bool = False,
    show_metric: str = "both",
) -> str:
    """
    Return a full HTML document with:
    - Hover tooltip on nsegment column (pure CSS :hover, not always-visible)
    - RdYlGn colour coding on % columns
    - ±CI95 column toggle
    - Lift column toggle
    """
    disp = tbl.copy()
    disp.index.name = None  # prevents pandas from rendering a second <tr> in <thead> for the index label
    orig_ids = disp.index.tolist()

    # Build column rename map
    rename = {
        "agg_bal":            "Bal%\u0394 (agg)",
        "agg_bal_ci":         "\u00b1CI Bal (agg)",
        "agg_acct":           "Acct%\u0394 (agg)",
        "agg_acct_ci":        "\u00b1CI Acct (agg)",
        "agg_lift_bal":       "Lift Bal (agg)",
        "agg_lift_bal_ci":    "\u00b1CI Lift Bal (agg)",
        "agg_lift_acct":      "Lift Acct (agg)",
        "agg_lift_acct_ci":   "\u00b1CI Lift Acct (agg)",
        "agg_n":              "N (agg)",
    }
    for c in ordered_comms:
        rename[f"{c}_bal"]          = f"{c} Bal%"
        rename[f"{c}_bal_ci"]       = f"{c} \u00b1CI"
        rename[f"{c}_acct"]         = f"{c} Acct%"
        rename[f"{c}_acct_ci"]      = f"{c} \u00b1CI Acct"
        rename[f"{c}_n"]            = f"{c} N"
        rename[f"{c}_lift_bal"]     = f"{c} Lift Bal"
        rename[f"{c}_lift_bal_ci"]  = f"{c} Lift CI±"
        rename[f"{c}_lift_acct"]    = f"{c} Lift Acct"
        rename[f"{c}_lift_acct_ci"] = f"{c} Lift Acct CI±"
    disp = disp.rename(columns=rename)

    # Drop columns based on toggles
    drop_cols = []
    # Always drop aggregate columns from table (still available in tbl for charts)
    _agg_keys = {"agg_bal", "agg_bal_ci", "agg_acct", "agg_acct_ci",
                 "agg_lift_bal", "agg_lift_bal_ci", "agg_lift_acct", "agg_lift_acct_ci", "agg_n"}
    drop_cols += [v for k, v in rename.items() if k in _agg_keys and v in disp.columns]
    if not show_n_cols:
        drop_cols += [v for k, v in rename.items() if k.endswith("_n") and v in disp.columns]
    # Drop raw-% CI from table (not lift CI — that's shown when show_lift=True)
    drop_cols += [v for k, v in rename.items()
                  if ("_ci" in k) and "_lift" not in k and v in disp.columns]
    # Lift is an exclusive mode: show lift columns OR raw % columns, never both
    if show_lift:
        # drop per-comm raw % columns; keep per-comm lift + lift CI columns
        drop_cols += [v for k, v in rename.items()
                      if v in disp.columns and "_ci" not in k and "_lift" not in k
                      and (k.endswith("_bal") or k.endswith("_acct"))]
    else:
        # drop per-comm lift columns (including lift CI); show raw % columns
        drop_cols += [v for k, v in rename.items() if ("_lift" in k) and v in disp.columns]
    # Metric filter: hide balance or accounts columns based on show_metric
    if show_metric == "balance":
        drop_cols += [v for k, v in rename.items() if "_acct" in k and v in disp.columns]
    elif show_metric == "accounts":
        drop_cols += [v for k, v in rename.items() if "_bal" in k and v in disp.columns]
    disp = disp.drop(columns=[c for c in drop_cols if c in disp.columns])

    pct_cols = [v for k, v in rename.items()
                if (k.endswith("_bal") or k.endswith("_acct") or k in ("agg_bal", "agg_acct", "agg_lift_bal", "agg_lift_acct") or "_lift" in k or "_ci" in k)
                and v in disp.columns]
    n_cols = [v for k, v in rename.items() if (k.endswith("_n") or k == "agg_n") and v in disp.columns]

    fmt = {col: _fmt_pct for col in pct_cols if col in disp.columns}
    fmt.update({col: "{:,.0f}" for col in n_cols if col in disp.columns})

    def _colour_col(col_series: pd.Series) -> pd.Series:
        return col_series.map(_rdylgn)

    # Colour lift columns with significance awareness (muted when CI crosses 0)
    lift_val_cols = [v for k, v in rename.items()
                     if ("_lift" in k) and "_ci" not in k and v in disp.columns]
    lift_ci_map   = {}  # display-name lift col → display-name CI col
    for k, v in rename.items():
        if "_lift" in k and "_ci" not in k and v in disp.columns:
            ci_key = k + "_ci"
            ci_v   = rename.get(ci_key, "")
            if ci_v in disp.columns:
                lift_ci_map[v] = ci_v

    plain_colour_cols = [v for k, v in rename.items()
                         if (k.endswith("_bal") or k.endswith("_acct")
                             or k in ("agg_bal", "agg_acct"))
                         and "_lift" not in k and "_ci" not in k
                         and v in disp.columns]

    styler = disp.style.format(fmt, na_rep="")
    if plain_colour_cols:
        styler = styler.apply(_colour_col, subset=plain_colour_cols, axis=0)
    # Plain coloring for all lift value columns
    if lift_val_cols:
        styler = styler.apply(_colour_col, subset=lift_val_cols, axis=0)
    if n_cols:
        styler = styler.apply(lambda s: s.map(_n_color), subset=n_cols, axis=0)

    table_html = styler.to_html()

    if show_lift and lift_ci_map:
        table_html = _inject_warn_flags(table_html, disp, list(lift_ci_map.items()))

    # Inject description tooltips on row-index cells
    if seg_labels or seg_desc:
        import re as _re
        def _inject_tooltip(m):
            tag_pre = m.group(1)
            cell_val = m.group(2)
            _id  = cell_val.strip()
            lbl  = (seg_labels or {}).get(_id, "")
            desc = (seg_desc   or {}).get(_id, "")
            tip  = f"{lbl} — {desc}" if lbl and desc else (lbl or desc)
            if not tip:
                return m.group(0)
            safe = tip.replace('"', "'")
            return f'{tag_pre} data-tip="{safe}">{cell_val}</th>'
        table_html = _re.sub(
            r'(<th[^>]*class="[^"]*row_heading[^"]*"[^>]*)>([^<]*)</th>',
            _inject_tooltip, table_html,
        )

    return f"""<!DOCTYPE html>
<html><head><meta charset="utf-8"><style>
  * {{ box-sizing: border-box; }}
  body {{ font-family: -apple-system, BlinkMacSystemFont, sans-serif;
          font-size: 12px; background: transparent; color: #fafafa; margin: 0; padding: 4px; }}
  table {{ border-collapse: collapse; width: 100%; }}
  thead th {{ background: #262730; color: #fafafa; padding: 6px 10px;
              text-align: left; position: sticky; top: 0; z-index: 2;
              border-bottom: 2px solid #555; white-space: nowrap; font-size: 11px; }}
  thead th:not(.blank) {{ cursor: pointer; user-select: none; }}
  thead th:not(.blank):hover {{ background: #3a3c4a; }}
  tbody td {{ padding: 4px 10px; border-bottom: 1px solid #333; white-space: nowrap; color: #111; }}
  tbody tr:hover td {{ outline: 1px solid #666; }}
  th.row_heading {{ background: #1c1e2a !important; font-size: 11px;
                    color: #aaa !important; font-weight: normal; cursor: help;
                    position: sticky; left: 0; z-index: 1;
                    min-width: 90px; padding: 4px 16px 4px 10px; }}
  th.blank {{ background: #262730 !important;
              position: sticky; left: 0; z-index: 3; min-width: 90px;
              cursor: pointer; user-select: none; }}
  th.blank:hover {{ background: #3a3c4a !important; }}
</style></head>
<body><div style="overflow:auto; max-height: 615px;">
{table_html}
</div><script>
(function(){{
  // Sort by data columns
  var ths = document.querySelectorAll('thead tr th:not(.blank)');
  var dir = {{}};
  ths.forEach(function(th, idx) {{
    th.addEventListener('click', function() {{
      var asc = !dir[idx]; dir[idx] = asc;
      var tbody = document.querySelector('tbody');
      var rows = Array.from(tbody.querySelectorAll('tr'));
      rows.sort(function(a, b) {{
        var tds_a = a.querySelectorAll('td');
        var tds_b = b.querySelectorAll('td');
        var av = (tds_a[idx] || {{innerText:''}}).innerText.replace(/[%,\u20ac]/g,'').trim();
        var bv = (tds_b[idx] || {{innerText:''}}).innerText.replace(/[%,\u20ac]/g,'').trim();
        var an = parseFloat(av), bn = parseFloat(bv);
        if (!isNaN(an) && !isNaN(bn)) return asc ? an-bn : bn-an;
        return asc ? av.localeCompare(bv) : bv.localeCompare(av);
      }});
      rows.forEach(function(r){{ tbody.appendChild(r); }});
    }});
  }});
  // Sort by index column (th.blank click)
  var blankTh = document.querySelector('thead tr th.blank');
  var idxDir = false;
  if (blankTh) {{
    blankTh.addEventListener('click', function() {{
      idxDir = !idxDir;
      var tbody = document.querySelector('tbody');
      var rows = Array.from(tbody.querySelectorAll('tr'));
      rows.sort(function(a, b) {{
        var av = (a.querySelector('th.row_heading') || {{innerText:''}}).innerText.trim();
        var bv = (b.querySelector('th.row_heading') || {{innerText:''}}).innerText.trim();
        var an = parseFloat(av), bn = parseFloat(bv);
        if (!isNaN(an) && !isNaN(bn)) return idxDir ? an-bn : bn-an;
        return idxDir ? av.localeCompare(bv) : bv.localeCompare(av);
      }});
      rows.forEach(function(r){{ tbody.appendChild(r); }});
    }});
  }}
}})();
(function(){{
  var tt=document.createElement('div');
  tt.style.cssText='position:fixed;background:#1e2030;color:#eee;font-size:11px;padding:5px 9px;border-radius:4px;border:1px solid #555;z-index:99999;pointer-events:none;display:none;max-width:340px;word-wrap:break-word;line-height:1.5;white-space:normal;box-shadow:0 2px 8px rgba(0,0,0,.5);';
  document.body.appendChild(tt);
  document.querySelectorAll('[data-tip]').forEach(function(el){{
    el.addEventListener('mouseenter',function(e){{tt.textContent=el.getAttribute('data-tip');tt.style.display='block';}});
    el.addEventListener('mousemove', function(e){{tt.style.left=(e.clientX+14)+'px';tt.style.top=(e.clientY+14)+'px';}});
    el.addEventListener('mouseleave',function(){{tt.style.display='none';}});
  }});
}})();
</script></body></html>"""


def _n_color(val: object) -> str:
    """Red-yellow-green background for N (sample size) cells."""
    try:
        n = float(val)  # type: ignore[arg-type]
        if n != n or n <= 0:  # NaN or zero
            return ""
        # 0 at n=30 → 1 at n=100; clamp outside
        ratio = min(1.0, max(0.0, (n - 30) / 70.0))
        if ratio < 0.5:
            t = ratio * 2
            r = int(248 + (255 - 248) * t)
            g = int(105 + (235 - 105) * t)
            b = int(107 + (132 - 107) * t)
        else:
            t = (ratio - 0.5) * 2
            r = int(255 + (99  - 255) * t)
            g = int(235 + (190 - 235) * t)
            b = int(132 + (123 - 132) * t)
        return f"background-color: rgb({r},{g},{b}); color: #000"
    except Exception:
        return ""


def _fmt_pct(v, na_str: str = "") -> str:
    """Format a fractional value (e.g. 0.05) as a readable % string."""
    if pd.isna(v):
        return na_str
    pct = v * 100
    if abs(pct) < 0.1:
        return f"{pct:.2f}%"
    if abs(pct) < 1.0:
        s = f"{pct:.1f}%"
        # guard: rounding can push 0.96% → "1.0%"; use 2dp instead
        return f"{pct:.2f}%" if s in ("1.0%", "-1.0%") else s
    return f"{pct:.0f}%"


def _inject_warn_flags(html: str, df: pd.DataFrame, warn_pairs: List[Tuple[str, str]]) -> str:
    """Inject ⚠️ into HTML table cells where CI > |lift| (statistically unreliable)."""
    import re as _re
    warn_cells: set = set()
    for lc, cc in warn_pairs:
        if lc in df.columns and cc in df.columns:
            for i in range(len(df)):
                lv, cv = df[lc].iat[i], df[cc].iat[i]
                if pd.notna(lv) and pd.notna(cv) and abs(cv) > abs(lv):
                    warn_cells.add((i, lc))
    if not warn_cells:
        return html
    col_names = list(df.columns)
    lift_positions = {lc: col_names.index(lc) for lc, _ in warn_pairs if lc in col_names}
    parts = html.split("<tbody>", 1)
    if len(parts) != 2:
        return html
    rows = parts[1].split("</tr>")
    row_idx = 0
    for ri in range(len(rows)):
        if "<tr" not in rows[ri]:
            continue
        tds = list(_re.finditer(r'(<td[^>]*>)(.*?)(</td>)', rows[ri]))
        for lc, col_pos in lift_positions.items():
            if (row_idx, lc) in warn_cells and col_pos < len(tds):
                m = tds[col_pos]
                rows[ri] = rows[ri].replace(m.group(0), f'{m.group(1)}{m.group(2)} ⚠️{m.group(3)}', 1)
        row_idx += 1
    return parts[0] + "<tbody>" + "</tr>".join(rows)


def _styled_html_table(
    styler,
    seg_labels: Optional[Dict[str, str]] = None,
    seg_desc: Optional[Dict[str, str]] = None,
    height: int = 400,
    warn_ci_pairs: Optional[List] = None,
) -> str:
    """Convert a Pandas Styler to a themed HTML doc with segment-ID hover tooltips."""
    import re as _re2
    html = styler.to_html()
    # Inject ⚠️ on lift cells where CI > |lift| (statistically unreliable)
    if warn_ci_pairs:
        html = _inject_warn_flags(html, styler.data, warn_ci_pairs)
    if seg_labels or seg_desc:
        def _tip(m):
            tag_pre  = m.group(1)
            cell_val = m.group(2)
            _id  = cell_val.strip()
            lbl  = (seg_labels or {}).get(_id, "")
            desc = (seg_desc   or {}).get(_id, "")
            tip  = f"{lbl} \u2014 {desc}" if lbl and desc else (lbl or desc)
            if not tip:
                return m.group(0)
            return f'{tag_pre} data-tip="{tip.replace(chr(34), chr(39))}">{cell_val}</th>'
        html = _re2.sub(
            r'(<th[^>]*class="[^"]*row_heading[^"]*"[^>]*)>([^<]*)</th>',
            _tip, html,
        )
    return f"""<!DOCTYPE html>
<html><head><meta charset="utf-8"><style>
  * {{ box-sizing: border-box; }}
  body {{ font-family: -apple-system, BlinkMacSystemFont, sans-serif;
          font-size: 12px; background: transparent; color: #fafafa; margin: 0; padding: 4px; }}
  table {{ border-collapse: collapse; width: 100%; }}
  thead th {{ background: #262730; color: #fafafa; padding: 6px 10px;
              text-align: left; position: sticky; top: 0; z-index: 2;
              border-bottom: 2px solid #555; white-space: nowrap; font-size: 11px; }}
  tbody td {{ padding: 4px 10px; border-bottom: 1px solid #333; white-space: nowrap; color: #111; }}
  tbody tr:hover td {{ outline: 1px solid #666; }}
  th.row_heading {{ background: #1c1e2a !important; font-size: 11px;
                    color: #ccc !important; font-weight: normal; cursor: help;
                    padding: 4px 6px 4px 6px; min-width: 50px; }}
  th.blank {{ background: #262730 !important; }}
</style></head><body>
<div style="overflow:auto; max-height:{height - 20}px;">{html}</div>
<script>(function(){{
  var tt=document.createElement('div');
  tt.style.cssText='position:fixed;background:#1e2030;color:#eee;font-size:11px;padding:5px 9px;border-radius:4px;border:1px solid #555;z-index:99999;pointer-events:none;display:none;max-width:340px;word-wrap:break-word;line-height:1.5;white-space:normal;box-shadow:0 2px 8px rgba(0,0,0,.5);';
  document.body.appendChild(tt);
  document.querySelectorAll('[data-tip]').forEach(function(el){{
    el.addEventListener('mouseenter',function(e){{tt.textContent=el.getAttribute('data-tip');tt.style.display='block';}});
    el.addEventListener('mousemove', function(e){{tt.style.left=(e.clientX+14)+'px';tt.style.top=(e.clientY+14)+'px';}});
    el.addEventListener('mouseleave',function(){{tt.style.display='none';}});
  }});
}})();
(function(){{
  // Truncate long text cells in the Description column with ellipsis + title tooltip
  var ths = document.querySelectorAll('thead th');
  var descIdx = -1;
  for (var i = 0; i < ths.length; i++) {{
    if (ths[i].innerText.trim() === 'Description') {{ descIdx = i - 1; break; }}
  }}
  if (descIdx < 0) return;
  document.querySelectorAll('tbody tr').forEach(function(tr) {{
    var tds = tr.querySelectorAll('td');
    if (descIdx < tds.length) {{
      var td = tds[descIdx];
      td.style.maxWidth = '260px';
      td.style.overflow = 'hidden';
      td.style.textOverflow = 'ellipsis';
      td.style.whiteSpace = 'nowrap';
      var text = td.innerText.trim();
      if (text) td.title = text;
    }}
  }});
}})();</script>
</body></html>"""


# ── Excel export ──────────────────────────────────────────────────────────────
def build_cc_bt_excel(tbl: pd.DataFrame, metric: str) -> bytes:
    """Build a colour-coded Excel for the CC Balance Transfer lift table.

    metric: 'conv' → conversion columns; 'amt' → amount columns.
    """
    buf = io.BytesIO()
    if metric == "conv":
        val_cols   = ["conv_treated", "conv_control", "conv_lift"]
        num_fmt    = "0.00%"
        ci_col     = "conv_lift_ci"
        col_labels = {
            "nsegment":     "Segment",
            "n_treated":    "N (treated)",
            "n_control":    "N (control)",
            "conv_treated": "Conv % (treated)",
            "conv_control": "Conv % (control)",
            "conv_lift":    "Conv Lift",
            "conv_lift_ci": "±CI",
        }
    else:
        val_cols   = ["amt_treated", "amt_control", "amt_lift"]
        num_fmt    = '$#,##0'
        ci_col     = "amt_lift_ci"
        col_labels = {
            "nsegment":    "Segment",
            "n_treated":   "N (treated)",
            "n_control":   "N (control)",
            "amt_treated": "Avg BT $ (treated)",
            "amt_control": "Avg BT $ (control)",
            "amt_lift":    "Amt Lift ($)",
            "amt_lift_ci": "±CI",
        }

    out = tbl.reset_index()
    ordered_cols = [c for c in col_labels if c in out.columns]
    out = out[ordered_cols].rename(columns=col_labels)

    with pd.ExcelWriter(buf, engine="xlsxwriter") as writer:
        out.to_excel(writer, sheet_name="CC BT Lift", index=False)
        wb = writer.book
        ws = writer.sheets["CC BT Lift"]

        hdr_fmt = wb.add_format({
            "bold": True, "bg_color": "#1C1E2A", "font_color": "#FFFFFF",
            "border": 1, "align": "center",
        })
        val_fmt  = wb.add_format({"num_format": num_fmt, "border": 1, "align": "center"})
        n_fmt    = wb.add_format({"num_format": "#,##0",  "border": 1, "align": "center"})
        seg_fmt  = wb.add_format({"border": 1, "bold": True})

        col_idx = {col: i for i, col in enumerate(out.columns)}

        # Header row styling
        for ci, h in enumerate(out.columns):
            ws.write(0, ci, h, hdr_fmt)

        # Column formats & widths
        lift_label = col_labels.get("conv_lift" if metric == "conv" else "amt_lift", "")
        for col_name, ci in col_idx.items():
            if col_name == "Segment":
                ws.set_column(ci, ci, 13, seg_fmt)
            elif col_name in ("N (treated)", "N (control)"):
                ws.set_column(ci, ci, 13, n_fmt)
            else:
                ws.set_column(ci, ci, 16, val_fmt)

        # 3-colour scale on all value columns
        nrows = len(out)
        for col_name in [col_labels.get(v, v) for v in val_cols]:
            if col_name not in col_idx:
                continue
            ci = col_idx[col_name]
            ws.conditional_format(
                1, ci, nrows, ci,
                {"type": "3_color_scale",
                 "min_color": "#F8696B", "mid_color": "#FFEB84", "max_color": "#63BE7B"},
            )

        # Data rows — explicit cell-by-cell write so formats stick
        for ri, row in enumerate(out.itertuples(index=False), start=1):
            for ci, col_name in enumerate(out.columns):
                v = row[ci]
                if col_name == "Segment":
                    ws.write(ri, ci, str(v) if pd.notna(v) else "", seg_fmt)
                elif col_name in ("N (treated)", "N (control)"):
                    ws.write(ri, ci, int(v) if pd.notna(v) else "", n_fmt)
                else:
                    ws.write(ri, ci, float(v) if pd.notna(v) else "", val_fmt)

    return buf.getvalue()


def build_excel(tbl: pd.DataFrame, ordered_comms: List[str]) -> bytes:
    buf = io.BytesIO()
    with pd.ExcelWriter(buf, engine="xlsxwriter") as writer:
        _drop = [c for c in tbl.columns if c.endswith("_ci") or c.startswith("agg_")]
        out = tbl.drop(columns=_drop).reset_index()
        out.to_excel(writer, sheet_name="Segments", index=False)
        wb = writer.book
        ws = writer.sheets["Segments"]
        pct_fmt  = wb.add_format({"num_format": "0.00%"})
        pct_keys = ["agg_bal", "agg_acct", "agg_lift_bal", "agg_lift_acct"] + \
                   [f"{c}_bal" for c in ordered_comms] + [f"{c}_acct" for c in ordered_comms]
        col_idx = {col: i for i, col in enumerate(out.columns)}
        for col in pct_keys:
            if col not in col_idx:
                continue
            ci = col_idx[col]
            ws.conditional_format(
                1, ci, len(out), ci,
                {"type": "3_color_scale", "min_color": "#F8696B",
                 "mid_color": "#FFEB84", "max_color": "#63BE7B"},
            )
            ws.set_column(ci, ci, 14, pct_fmt)
        ws.set_column(0, 0, 12)
    return buf.getvalue()


def _build_sim_excel(
    dl_detail: pd.DataFrame,
    summary_df: pd.DataFrame,
    ordered_comms: List[str],
    metric_lbl: str,   # e.g. " Bal%" or " Acct%"
    include_lift: bool = True,
    include_n: bool = True,
    include_prop: bool = False,
    include_summary: bool = True,
    selected_comms: Optional[List[str]] = None,   # None = all comms
) -> bytes:
    """Build a richly formatted Excel for the Audience Simulator download."""
    buf = io.BytesIO()

    # Apply communication filter
    _comms = [c for c in ordered_comms if selected_comms is None or c in selected_comms]

    # Build export df: fixed cols + interleaved [Lift%, N, Prop%] per comm
    bk = dl_detail.copy()

    export_cols_ordered: List[str] = [
        c for c in ["Role", "Label", "Description"] if c in bk.columns
    ]
    for _c in _comms:
        _lc = f"{_c}{metric_lbl}"
        _nc = f"{_c} N"
        _pc = f"{_c} Prop%"
        if include_lift and _lc in bk.columns:
            export_cols_ordered.append(_lc)
        if _nc in bk.columns:
            _total = bk[_nc].sum()
            bk[_pc] = (bk[_nc] / _total) if _total > 0 else np.nan
            if include_n:
                export_cols_ordered.append(_nc)
            if include_prop:
                export_cols_ordered.append(_pc)
    bk = bk[[c for c in export_cols_ordered if c in bk.columns]]

    with pd.ExcelWriter(buf, engine="xlsxwriter") as _ew:
        wb = _ew.book

        # ── Shared formats ────────────────────────────────────────────────────
        _hdr = wb.add_format({
            "bold": True, "bg_color": "#1C1E2A", "font_color": "#FFFFFF",
            "border": 1, "border_color": "#000000",
            "align": "center", "valign": "vcenter", "text_wrap": True,
        })
        _idx = wb.add_format({
            "bg_color": "#1C1E2A", "font_color": "#AAAAAA",
            "border": 1, "border_color": "#000000", "font_size": 10,
        })
        _pct   = wb.add_format({"num_format": "0%",     "border": 1, "align": "center"})
        _nfmt  = wb.add_format({"num_format": "#,##0",  "border": 1, "align": "center"})
        _prop  = wb.add_format({"num_format": "0.0%",   "border": 1, "align": "center"})
        _txt   = wb.add_format({"border": 1, "text_wrap": True, "valign": "top"})
        _eur   = wb.add_format({"num_format": "\u20ac#,##0", "border": 1, "align": "right"})
        _role  = {
            "NOT": wb.add_format({"bg_color": "#FFDDDD", "font_color": "#880000",
                                   "border": 1, "bold": True, "align": "center"}),
            "AND": wb.add_format({"bg_color": "#E8D5FF", "font_color": "#5500AA",
                                   "border": 1, "bold": True, "align": "center"}),
            "OR":  wb.add_format({"bg_color": "#D0E8FF", "font_color": "#003399",
                                   "border": 1, "bold": True, "align": "center"}),
        }

        # ── Breakdown sheet ───────────────────────────────────────────────────
        ws = wb.add_worksheet("Breakdown")
        _ew.sheets["Breakdown"] = ws
        all_cols = list(bk.columns)
        nrows = len(bk)

        # Header row
        ws.set_row(0, 30)
        ws.write(0, 0, "Segment", _hdr)
        for _ci, _col in enumerate(all_cols, start=1):
            ws.write(0, _ci, _col, _hdr)

        # Data rows
        for _ri, (_seg, _row) in enumerate(bk.iterrows(), start=1):
            ws.write(_ri, 0, str(_seg), _idx)
            for _ci, _col in enumerate(all_cols, start=1):
                _v = _row[_col]
                if _col == "Role":
                    ws.write(_ri, _ci, _v if pd.notna(_v) else "", _role.get(str(_v), _txt))
                elif _col.endswith(metric_lbl):           # lift % col
                    ws.write(_ri, _ci, float(_v) if pd.notna(_v) else "", _pct)
                elif _col.endswith(" N"):
                    ws.write(_ri, _ci, int(_v) if pd.notna(_v) else "", _nfmt)
                elif _col.endswith("Prop%"):
                    ws.write(_ri, _ci, float(_v) if pd.notna(_v) else "", _prop)
                else:
                    ws.write(_ri, _ci, str(_v) if pd.notna(_v) else "", _txt)

        # Conditional colour scales
        if nrows > 0:
            for _ci, _col in enumerate(all_cols, start=1):
                if _col.endswith(metric_lbl):
                    ws.conditional_format(1, _ci, nrows, _ci, {
                        "type": "3_color_scale",
                        "min_color": "#F8696B", "mid_color": "#FFEB84", "max_color": "#63BE7B",
                        "min_type": "percentile", "min_value": 10,
                        "mid_type": "percentile", "mid_value": 50,
                        "max_type": "percentile", "max_value": 90,
                    })
                elif _col.endswith(" N"):
                    ws.conditional_format(1, _ci, nrows, _ci, {
                        "type": "3_color_scale",
                        "min_color": "#F8696B", "mid_color": "#FFEB84", "max_color": "#63BE7B",
                        "min_type": "num", "min_value": 30,
                        "mid_type": "num", "mid_value": 65,
                        "max_type": "num", "max_value": 100,
                    })
                elif _col.endswith("Prop%"):
                    ws.conditional_format(1, _ci, nrows, _ci, {
                        "type": "3_color_scale",
                        "min_color": "#F8696B", "mid_color": "#FFEB84", "max_color": "#63BE7B",
                    })

        # Column widths
        ws.set_column(0, 0, 13)  # segment index
        _col_w = {"Role": 8, "Label": 22, "Description": 42}
        for _ci, _col in enumerate(all_cols, start=1):
            _w = _col_w.get(
                _col,
                10 if _col.endswith(" N") or _col.endswith("Prop%")
                else 13 if _col.endswith(metric_lbl)
                else 14,
            )
            ws.set_column(_ci, _ci, _w)

        ws.freeze_panes(1, 1)
        ws.autofilter(0, 0, nrows, len(all_cols))

        # ── Summary sheet ─────────────────────────────────────────────────────
        if include_summary and len(summary_df) > 0:
            ws2 = wb.add_worksheet("Summary")
            _ew.sheets["Summary"] = ws2
            _sum_cols = [c for c in ["Communication", "Expected Lift", "Comm Users",
                                      "Projected Bal \u20ac", "Proj. Accounts"] if c in summary_df.columns]
            ws2.set_row(0, 30)
            for _ci, _col in enumerate(_sum_cols):
                ws2.write(0, _ci, _col, _hdr)
            _comm_fmt = wb.add_format({"bold": True, "border": 1})
            for _ri, _row in summary_df[_sum_cols].iterrows():
                for _ci, _col in enumerate(_sum_cols):
                    _v = _row[_col]
                    if _col == "Communication":
                        ws2.write(_ri + 1, _ci, _v, _comm_fmt)
                    elif _col == "Expected Lift":
                        ws2.write(_ri + 1, _ci, float(_v) if pd.notna(_v) else "", _pct)
                    elif _col in ("Comm Users", "Proj. Accounts"):
                        ws2.write(_ri + 1, _ci, int(_v) if pd.notna(_v) else "", _nfmt)
                    elif _col == "Projected Bal \u20ac":
                        ws2.write(_ri + 1, _ci, float(_v) if pd.notna(_v) else "", _eur)
            if len(summary_df) > 0 and "Expected Lift" in _sum_cols:
                _lci = _sum_cols.index("Expected Lift")
                ws2.conditional_format(1, _lci, len(summary_df), _lci, {
                    "type": "3_color_scale",
                    "min_color": "#F8696B", "mid_color": "#FFEB84", "max_color": "#63BE7B",
                })
            _w2 = {"Communication": 16, "Expected Lift": 14, "Comm Users": 14,
                   "Projected Bal \u20ac": 20, "Proj. Accounts": 16}
            for _ci, _col in enumerate(_sum_cols):
                ws2.set_column(_ci, _ci, _w2.get(_col, 14))
            ws2.freeze_panes(1, 0)

    buf.seek(0)
    return buf.read()


# ── PPT export ────────────────────────────────────────────────────────────────
def build_pptx(tbl: pd.DataFrame, ordered_comms: List[str]) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Title slide
    sl = prs.slides.add_slide(blank_layout)
    txb = sl.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
    tf = txb.text_frame
    tf.text = "Affinity Explorer"
    tf.paragraphs[0].runs[0].font.size = Pt(36)
    tf.paragraphs[0].runs[0].font.bold = True
    tf.paragraphs[0].runs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = RGBColor(0x1C, 0x1E, 0x2A)

    # Table slide (top 20 rows, key cols)
    key_cols = ["agg_bal", "agg_lift_bal", "agg_acct", "agg_n"]
    key_cols = [c for c in key_cols if c in tbl.columns]
    top20 = tbl[key_cols].dropna(subset=["agg_bal"]).sort_values("agg_bal", ascending=False).head(20)

    sl2 = prs.slides.add_slide(blank_layout)
    sl2.background.fill.solid()
    sl2.background.fill.fore_color.rgb = RGBColor(0x1C, 0x1E, 0x2A)
    txb2 = sl2.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(12), Inches(0.6))
    txb2.text_frame.text = "Top 20 Segments — Balance % Change"
    txb2.text_frame.paragraphs[0].runs[0].font.size = Pt(18)
    txb2.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    rows_n = len(top20) + 1
    cols_n = len(key_cols) + 1
    tbl_shape = sl2.shapes.add_table(rows_n, cols_n, Inches(0.3), Inches(0.9), Inches(12.7), Inches(6.2))
    ptbl = tbl_shape.table

    hdr_fill = RGBColor(0x26, 0x27, 0x30)
    _pptx_rename = {"agg_bal": "Balance %", "agg_lift_bal": "Lift (Balance)", "agg_acct": "Accounts %", "agg_n": "Customers"}
    headers = ["Segment"] + [_pptx_rename.get(c, c.replace("_", " ").title()) for c in key_cols]
    for ci, h in enumerate(headers):
        cell = ptbl.cell(0, ci)
        cell.text = h
        cell.text_frame.paragraphs[0].runs[0].font.size = Pt(9)
        cell.text_frame.paragraphs[0].runs[0].font.bold = True
        cell.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        cell.fill.solid()
        cell.fill.fore_color.rgb = hdr_fill

    for ri, (seg, row) in enumerate(top20.iterrows(), start=1):
        ptbl.cell(ri, 0).text = str(seg)
        ptbl.cell(ri, 0).text_frame.paragraphs[0].runs[0].font.size = Pt(8)
        for ci, col in enumerate(key_cols, start=1):
            val = row[col]
            if pd.isna(val):
                ptbl.cell(ri, ci).text = ""
            elif col.endswith("_n"):
                ptbl.cell(ri, ci).text = f"{int(val):,}"
            else:
                ptbl.cell(ri, ci).text = f"{val:.2%}"
            ptbl.cell(ri, ci).text_frame.paragraphs[0].runs[0].font.size = Pt(8)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ── Excel export dialog ───────────────────────────────────────────────────────
@st.dialog("Export to Excel")
def _sim_excel_dialog(
    dl_detail: pd.DataFrame,
    summary_df: pd.DataFrame,
    ordered_comms: list,
    metric_lbl: str,
) -> None:
    """Modal dialog: choose which communications and column types to export."""
    st.markdown("**Communications to include:**")
    _sel_comms = [
        c for c in ordered_comms
        if st.checkbox(c, value=True, key=f"_xldlg_{c}")
    ]
    if not _sel_comms:
        st.warning("Select at least one communication.")
        return

    st.markdown("---")
    st.markdown("**Column types:**")
    _c1, _c2, _c3 = st.columns(3)
    inc_lift    = _c1.checkbox("Lift%",       value=True,  key="_xldlg_lift")
    inc_n       = _c2.checkbox("N (treated)", value=True,  key="_xldlg_n")
    inc_prop    = _c3.checkbox("Pop %",       value=False, key="_xldlg_prop")
    inc_summary = st.checkbox("Include Summary sheet", value=True, key="_xldlg_sum")

    if not inc_lift and not inc_n and not inc_prop:
        st.warning("Select at least one column type.")
        return

    _bytes = _build_sim_excel(
        dl_detail, summary_df, ordered_comms, metric_lbl,
        include_lift=inc_lift,
        include_n=inc_n,
        include_prop=inc_prop,
        include_summary=inc_summary,
        selected_comms=_sel_comms,
    )
    st.download_button(
        "Download",
        _bytes,
        file_name="targeting_selection.xlsx",
        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        use_container_width=True,
    )


# ── Page config ───────────────────────────────────────────────────────────────
st.set_page_config(page_title="Affinity Explorer", layout="wide")

st.markdown("""
<style>
span[data-baseweb="tag"] { background-color: #31333F !important; }
span[data-baseweb="tag"] span { color: #FAFAFA !important; }
span[data-baseweb="tag"] button { color: #FAFAFA !important; }
div.block-container { padding-top: 3rem !important; }
</style>
""", unsafe_allow_html=True)
_title_col, _camp_col, _ = st.columns([3, 2, 3])
_title_col.title("Affinity Explorer")
with _camp_col:
    active_campaign = st.selectbox(
        "Campaign",
        options=list(CAMPAIGNS.keys()),
        format_func=lambda k: CAMPAIGNS[k]["label"],
        key="active_campaign",
        label_visibility="collapsed",
    )

# ── Load audience profile (shared across all campaigns) ──────────────────────
_aud_profile_path = _APP_DIR / "audience_profile.csv"
try:
    with open(_aud_profile_path, "rb") as _f:
        _aud_df = _load_audience_profile(_f.read())
except FileNotFoundError:
    _aud_df = None

# ── shared multi-table HTML renderer used by all campaign branches ───────────
def _two_tables_html(items: list, height: int) -> str:
    """Render multiple (label, styler) pairs into one HTML doc — zero spacing between tables."""
    import re as _re3
    _css = (
        "* { box-sizing: border-box; }"
        "body { font-family: -apple-system, BlinkMacSystemFont, sans-serif;"
        "       font-size: 12px; background: transparent; color: #fafafa; margin: 0; padding: 4px; }"
        "table { border-collapse: collapse; width: 100%; table-layout: fixed; }"
        "thead th { background: #262730; color: #fafafa; padding: 6px 10px;"
        "           text-align: left; position: sticky; top: 0; z-index: 2;"
        "           border-bottom: 2px solid #555; white-space: nowrap; font-size: 11px; }"
        "tbody td { padding: 4px 10px; border-bottom: 1px solid #333; white-space: nowrap; color: #111; }"
        "tbody td:first-child { white-space: normal; word-break: break-word; }"
        "tbody tr:hover td { outline: 1px solid #666; }"
        "th.row_heading { background: #1c1e2a !important; font-size: 11px;"
        "                 color: #ccc !important; font-weight: normal; cursor: help;"
        "                 padding: 4px 6px 4px 6px; min-width: 50px; }"
        "th.blank { background: #262730 !important; }"
        ".tbl-lbl { font-size: 11px; font-weight: 600; color: #bbb;"
        "           padding: 6px 2px 2px 2px; margin-top: 2px; }"
        ".tbl-lbl:first-child { margin-top: 0; padding-top: 2px; }"
    )
    _blocks = []
    for _item in items:
        _lbl = _item[0]; _styler = _item[1]
        _warn_p = _item[2] if len(_item) > 2 else None
        _h = _styler.to_html()
        if _warn_p:
            _h = _inject_warn_flags(_h, _styler.data, _warn_p)
        def _tip(m):
            _tp = m.group(1); _cv = m.group(2); _id = _cv.strip()
            _l = (SEGMENT_LABELS or {}).get(_id, ""); _d2 = (SEGMENT_DESCRIPTIONS or {}).get(_id, "")
            _t = f"{_l} \u2014 {_d2}" if _l and _d2 else (_l or _d2)
            if not _t: return m.group(0)
            return f'{_tp} data-tip="{_t.replace(chr(34), chr(39))}">{_cv}</th>'
        _h = _re3.sub(r'(<th[^>]*class="[^"]*row_heading[^"]*"[^>]*)>([^<]*)</th>', _tip, _h)
        _ncols_inj = 1 + len(_styler.data.columns)
        _pct_map = {5: [10, 46, 16, 18, 10], 4: [14, 42, 25, 19], 6: [10, 30, 16, 18, 16, 10]}
        _col_pcts = _pct_map.get(_ncols_inj, [round(100 / _ncols_inj)] * _ncols_inj)
        _cg = '<colgroup>' + ''.join(f'<col style="width:{p}%">' for p in _col_pcts) + '</colgroup>'
        _h = _re3.sub(r'(<table[^>]*>)', r'\1' + _cg, _h, count=1)
        _blocks.append(f'<div class="tbl-lbl">{_lbl}</div>{_h}')
    _body = "".join(_blocks)
    _js = (
        "(function(){var tt=document.createElement('div');"
        "tt.style.cssText='position:fixed;background:#1e2030;color:#eee;font-size:11px;"
        "padding:5px 9px;border-radius:4px;border:1px solid #555;z-index:99999;"
        "pointer-events:none;display:none;max-width:340px;word-wrap:break-word;"
        "line-height:1.5;white-space:normal;box-shadow:0 2px 8px rgba(0,0,0,.5);';"
        "document.body.appendChild(tt);"
        "document.querySelectorAll('[data-tip]').forEach(function(el){"
        "el.addEventListener('mouseenter',function(){tt.textContent=el.getAttribute('data-tip');tt.style.display='block';});"
        "el.addEventListener('mousemove',function(e){tt.style.left=(e.clientX+14)+'px';tt.style.top=(e.clientY+14)+'px';});"
        "el.addEventListener('mouseleave',function(){tt.style.display='none';});"
        "})})()"
    )
    return (
        f"<!DOCTYPE html><html><head><meta charset='utf-8'>"
        f"<style>{_css}</style></head><body>"
        f"<div style='overflow:auto;max-height:{height-20}px;'>{_body}</div>"
        f"<script>{_js}</script></body></html>"
    )


# ── shared multi-table HTML renderer used by all campaign branches ─────────────
def _two_tables_html(items: list, height: int) -> str:
    """Render multiple (label, styler) pairs into one HTML doc — zero spacing between tables."""
    import re as _re3
    _css = (
        "* { box-sizing: border-box; }"
        "body { font-family: -apple-system, BlinkMacSystemFont, sans-serif;"
        "       font-size: 12px; background: transparent; color: #fafafa; margin: 0; padding: 4px; }"
        "table { border-collapse: collapse; width: 100%; table-layout: fixed; }"
        "thead th { background: #262730; color: #fafafa; padding: 6px 10px;"
        "           text-align: left; position: sticky; top: 0; z-index: 2;"
        "           border-bottom: 2px solid #555; white-space: nowrap; font-size: 11px; }"
        "tbody td { padding: 4px 10px; border-bottom: 1px solid #333; white-space: nowrap; color: #111; }"
        "tbody td:first-child { white-space: normal; word-break: break-word; }"
        "tbody tr:hover td { outline: 1px solid #666; }"
        "th.row_heading { background: #1c1e2a !important; font-size: 11px;"
        "                 color: #ccc !important; font-weight: normal; cursor: help;"
        "                 padding: 4px 6px 4px 6px; min-width: 50px; }"
        "th.blank { background: #262730 !important; }"
        ".tbl-lbl { font-size: 11px; font-weight: 600; color: #bbb;"
        "           padding: 6px 2px 2px 2px; margin-top: 2px; }"
        ".tbl-lbl:first-child { margin-top: 0; padding-top: 2px; }"
    )
    _blocks = []
    for _item in items:
        _lbl = _item[0]; _styler = _item[1]
        _warn_p = _item[2] if len(_item) > 2 else None
        _h = _styler.to_html()
        if _warn_p:
            _h = _inject_warn_flags(_h, _styler.data, _warn_p)
        def _tip(m):
            _tp = m.group(1); _cv = m.group(2); _id = _cv.strip()
            _l = (SEGMENT_LABELS or {}).get(_id, ""); _d2 = (SEGMENT_DESCRIPTIONS or {}).get(_id, "")
            _t = f"{_l} \u2014 {_d2}" if _l and _d2 else (_l or _d2)
            if not _t: return m.group(0)
            return f'{_tp} data-tip="{_t.replace(chr(34), chr(39))}">{_cv}</th>'
        _h = _re3.sub(r'(<th[^>]*class="[^"]*row_heading[^"]*"[^>]*)>([^<]*)</th>', _tip, _h)
        _ncols_inj = 1 + len(_styler.data.columns)
        _pct_map = {5: [10, 46, 16, 18, 10], 4: [14, 42, 25, 19], 6: [10, 30, 16, 18, 16, 10]}
        _col_pcts = _pct_map.get(_ncols_inj, [round(100 / _ncols_inj)] * _ncols_inj)
        _cg = '<colgroup>' + ''.join(f'<col style="width:{p}%">' for p in _col_pcts) + '</colgroup>'
        _h = _re3.sub(r'(<table[^>]*>)', r'\1' + _cg, _h, count=1)
        _blocks.append(f'<div class="tbl-lbl">{_lbl}</div>{_h}')
    _body = "".join(_blocks)
    _js = (
        "(function(){var tt=document.createElement('div');"
        "tt.style.cssText='position:fixed;background:#1e2030;color:#eee;font-size:11px;"
        "padding:5px 9px;border-radius:4px;border:1px solid #555;z-index:99999;"
        "pointer-events:none;display:none;max-width:340px;word-wrap:break-word;"
        "line-height:1.5;white-space:normal;box-shadow:0 2px 8px rgba(0,0,0,.5);';"
        "document.body.appendChild(tt);"
        "document.querySelectorAll('[data-tip]').forEach(function(el){"
        "el.addEventListener('mouseenter',function(){tt.textContent=el.getAttribute('data-tip');tt.style.display='block';});"
        "el.addEventListener('mousemove',function(e){tt.style.left=(e.clientX+14)+'px';tt.style.top=(e.clientY+14)+'px';});"
        "el.addEventListener('mouseleave',function(){tt.style.display='none';});"
        "})})()"
    )
    return (
        f"<!DOCTYPE html><html><head><meta charset='utf-8'>"
        f"<style>{_css}</style></head><body>"
        f"<div style='overflow:auto;max-height:{height-20}px;'>{_body}</div>"
        f"<script>{_js}</script></body></html>"
    )


# ══════════════════════════════════════════════════════════
# SELECTCHK CAMPAIGN
# ══════════════════════════════════════════════════════════
if active_campaign == "SELECTCHK":

    # ── Load raw CSV ──────────────────────────────────────────────────────────
    try:
        with open(_APP_DIR / "SELECTCHK_campaign.csv", "rb") as f:
            raw_bytes = f.read()
        df_raw = _load_csv(raw_bytes)
    except FileNotFoundError:
        st.error("SELECTCHK_campaign.csv not found — please upload or provide a SELECTCHK_campaign.csv file.")
        st.stop()

    missing_cols = [c for c in REQUIRED_COLS if c not in df_raw.columns]
    if missing_cols:
        st.error("Missing required columns: " + ", ".join(missing_cols))
        st.stop()

    # Max unique users in the dataset — used as the upper bound for the min-N slider
    _max_slider_n = max(30, int(df_raw["alpha_key"].nunique()))

    # ── Sidebar ───────────────────────────────────────────────────────────────────
    with st.sidebar:
        # ── Segment lookup ────────────────────────────────────────────────────
        _seg_ids = _all_segment_ids(df_raw)
        _seg_lookup_widget(_seg_ids, SEGMENT_LABELS, SEGMENT_DESCRIPTIONS)

        st.subheader("Date range")
        _dates = pd.to_datetime(df_raw["start_date"], errors="coerce").dropna()
        _d_min = _dates.min().date()
        _d_max = _dates.max().date()
        _dcol1, _dcol2 = st.columns(2)
        with _dcol1:
            date_from = st.date_input("From", value=_d_min, min_value=_d_min, max_value=_d_max)
        with _dcol2:
            date_to = st.date_input("To", value=_d_max, min_value=_d_min, max_value=_d_max)

        st.divider()
        with st.expander("Advanced settings", expanded=False):
            all_mode = False
            min_n = _MIN_N_DEFAULT  # 30 — segments with fewer treated customers are hidden
            bal_clip_pct = 5  # clip below 5th and above 95th percentile (hardcoded)
            bal_baseline_min = 25.0  # exclude segments whose avg starting balance is below €25
            st.caption(
                "**Minimum N** — segments with fewer than 30 treated customers are hidden from "
                "the table. Too few observations make the mean unreliable and the CI very wide. "
                "Value: **30 customers**."
            )
            st.caption(
                "**Outlier clipping** — the bottom 5% and top 5% of `balance_pct_change` values "
                "are removed *before* computing any means or lifts. This prevents a handful of "
                "extreme balance swings (e.g. account closures, large one-off transfers) from "
                "distorting the segment averages."
            )
            st.caption(
                "**Baseline balance filter** — segments whose average *starting* balance across "
                "treated customers is below **$25** are excluded. These are near-zero balance "
                "accounts where % change is noisy and economically meaningless."
            )

        # Show columns — metric toggle lives inside the Data tab (see tab_table section)
        show_n_cols     = False
        _show_metric_val = "both"  # default; overridden by widget inside tab_table


    # ── Preprocess ────────────────────────────────────────────────────────────────
    df = preprocess(
        df_raw,
        date_min=str(date_from),
        date_max=str(date_to),
        bal_clip_pct=float(bal_clip_pct),
    )

    # ── ordered_comms from session state (checkboxes live inside Table tab) ─────────
    _all_present_comms = [c for c in COMM_ORDER if c in df["communication"].unique()]
    ordered_comms = [c for c in _all_present_comms if st.session_state.get(f"cb_{c}", True)]
    if not ordered_comms:
        ordered_comms = _all_present_comms[:]

    # ── Tabs ──────────────────────────────────────────────────────────────────────
    tab_explorer, tab_table, tab_simulator, tab_audit = st.tabs([
        "Segment Explorer",
        "Data",
        "Audience Simulator",
        "Data Quality",
    ])


    # ══════════════════════════════════════════════════════════
    # SEGMENT EXPLORER TAB
    # ══════════════════════════════════════════════════════════
    with tab_explorer:
        st.subheader("Segment catalogue")
        st.caption("Browse all segments by category. Unique Customers = contacted customers in the active date range.")
        if not SEGMENT_LABELS:
            st.info("No segment descriptions found. Ensure segment_descriptions.csv is present in the app directory.")
        else:
            # ── Build base table from CSV (all segments defined) ──────────────────
            _exp_base = pd.DataFrame({
                "Segment ID": list(SEGMENT_LABELS.keys()),
                "Group":      list(SEGMENT_LABELS.values()),
                "Description": [SEGMENT_DESCRIPTIONS.get(k, "") for k in SEGMENT_LABELS],
            })
            _user_counts = (
                df[df["contact_flag"] == 1]
                .groupby("nsegment")["alpha_key"]
                .nunique()
                .rename("Unique Customers")
                .reset_index()
                .rename(columns={"nsegment": "Segment ID"})
            )
            _exp_base = _exp_base.merge(_user_counts, on="Segment ID", how="left")
            _exp_base["Unique Customers"] = _exp_base["Unique Customers"].fillna(0).astype(int)
            _exp_base = _exp_base.sort_values("Segment ID").reset_index(drop=True)

            # ── Full segment table (always visible) ────────────────────────────────
            _exp_styler = (
                _exp_base.set_index("Segment ID").rename_axis(None)
                .style
                .apply(lambda s: s.map(_n_color), subset=["Unique Customers"], axis=0)
                .format(na_rep="")
            )
            _exp_h = max(300, min(700, 60 + len(_exp_base) * 26))
            components.html(
                _styled_html_table(_exp_styler, SEGMENT_LABELS, SEGMENT_DESCRIPTIONS, height=_exp_h),
                height=_exp_h, scrolling=True,
            )

            st.divider()

            # ── Group size chart ───────────────────────────────────────────────────────
            _grp_sz = (
                _exp_base.groupby("Group")
                .agg(Segments=("Segment ID", "count"), Customers=("Unique Customers", "sum"))
                .reset_index()
                .sort_values("Customers", ascending=False)
            )
            _fig_gsz = px.bar(
                _grp_sz, x="Group", y="Segments",
                color="Segments", color_continuous_scale="Blues",
                text="Segments",
                title="Segments per group",
            )
            _fig_gsz.update_xaxes(tickangle=-30)
            _fig_gsz.update_traces(textposition="outside")
            _fig_gsz.update_layout(
                coloraxis_showscale=False, height=360,
                yaxis_range=[0, int(_grp_sz["Segments"].max()) * 1.20],
            )
            st.plotly_chart(_fig_gsz, width='stretch')

            _fig_gus = px.bar(
                _grp_sz, x="Group", y="Customers",
                color="Customers", color_continuous_scale="Teal",
                text=_grp_sz["Customers"].map(lambda v: f"{v:,}"),
                title="Unique customers per group (active date range)",
            )
            _fig_gus.update_xaxes(tickangle=-30)
            _fig_gus.update_traces(textposition="outside")
            _fig_gus.update_layout(
                coloraxis_showscale=False, height=360,
                yaxis_range=[0, int(_grp_sz["Customers"].max()) * 1.20],
            )
            st.plotly_chart(_fig_gus, width='stretch')


    # ══════════════════════════════════════════════════════════
    # TABLE TAB
    # ══════════════════════════════════════════════════════════
    with tab_table:
        # ── Metric toggle ─────────────────────────────────────────────────────────
        _metric_opts = ["Balance only", "Balance & Accounts", "Accounts only"]
        _show_metric_radio = st.radio(
            "Show columns",
            _metric_opts,
            horizontal=True,
            key="show_metric_radio",
            label_visibility="visible",
        )
        _show_metric_map = {"Balance & Accounts": "both", "Balance only": "balance", "Accounts only": "accounts"}
        _show_metric_val = _show_metric_map[_show_metric_radio]

        # ── Row 1: comm toggle strip + Show N + Show lift ─────────────────────────
        _n_pc = len(_all_present_comms)
        _row1 = st.columns(_n_pc + 2)
        for _i, _c in enumerate(_all_present_comms):
            _row1[_i].checkbox(_c, value=False, key=f"cb_{_c}")
        show_n_cols = _row1[_n_pc].checkbox("Show N", value=False, key="show_n_cols")
        show_lift = _row1[_n_pc + 1].checkbox(
            "Show lift",
            value=True,
            key="tbl_show_lift",
            help="Lift = treatment group mean − control group mean in the same segment. "
                 "Positive = the communication added value beyond background trends. "
                 "Enable this to see whether the % change is *caused by* the communication, not just correlated.",
        )

        # ── Combo mode vs standard mode ──────────────────────────────────────────
        _checked_comms = [c for c in _all_present_comms if st.session_state.get(f"cb_{c}", False)]
        _combo_mode = len(_checked_comms) > 0
        if _combo_mode:
            _per_comm_t = {c: set(df[(df["contact_flag"] == 1) & (df["communication"] == c)]["alpha_key"]) for c in _checked_comms}
            _combo_treated_keys = frozenset(set.intersection(*_per_comm_t.values())) if _per_comm_t else frozenset()
            ordered_comms = _checked_comms
            st.caption(
                f"Combo mode — treatment: **{len(_combo_treated_keys):,}** customers who received all of: "
                f"{', '.join(_checked_comms)}. Control: `contact_flag=0` per communication."
            )
        else:
            _combo_treated_keys = None
            ordered_comms = _all_present_comms[:]
            pass

        with st.spinner("Computing..."):
            tbl = build_table(df, ordered_comms, all_mode, min_n, bal_baseline_min,
                              combo_treated_keys=_combo_treated_keys, _ver=2)

        if tbl.empty:
            # Show diagnostic info so user knows exactly why nothing appears
            _n_treated  = int(df[df["contact_flag"] == 1]["alpha_key"].nunique())
            _n_ctrl     = int(df[df["control_flag"] == 1]["alpha_key"].nunique())
            _comms_seen = sorted(df["communication"].unique().tolist())
            _segs_seen  = int(df["nsegment"].nunique())
            st.warning(
                f"**No segments pass the current filters.**  \n"
                f"Treated: **{_n_treated:,}** | Control: **{_n_ctrl:,}** | "
                f"Segments: **{_segs_seen:,}** | Communications: **{', '.join(_comms_seen)}**  \n\n"
                f"Try widening the date range, or check the data has treated and control customers."
            )
        else:
            html_tbl = style_tbl(
                tbl, ordered_comms,
                seg_labels=SEGMENT_LABELS,
                seg_desc=SEGMENT_DESCRIPTIONS,
                show_n_cols=show_n_cols,
                show_lift=show_lift,
                show_metric=_show_metric_val,
            )
            _tbl_h = max(300, min(680, 55 + len(tbl) * 32))
            components.html(html_tbl, height=_tbl_h, scrolling=True)

            # ── Downloads ────────────────────────────────────────────────────
            _, _sdl_col, _ = st.columns([2, 3, 2])
            _sdl_col.download_button(
                label="Download Excel",
                data=build_excel(tbl, ordered_comms),
                file_name="selectchk_segment_lift.xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                use_container_width=True,
                key="selectchk_dl_xlsx",
            )

            if show_lift:
                st.caption(
                    "**Blank cells** = fewer than 30 customers in that segment × communication (too small to trust)."
                )

            if show_lift:
                with st.expander("How to read Lift \u00b1 CI", expanded=False):
                    st.markdown(
                        "<p style='font-size:1.05rem;line-height:1.7'>"
                        "<b>Segment A</b> — Lift 5.0% \u00b1 2.0% → range 3.0%–7.0%. "
                        "The range stays <b>above zero</b> → the lift is <b>statistically reliable</b>.<br><br>"
                        "<b>Segment B</b> — Lift 2.0% \u00b1 3.5% → range −1.5%–5.5%. "
                        "The range <b>crosses zero</b> → the lift could actually be negative; "
                        "we can't be confident it's real."
                        "</p>",
                        unsafe_allow_html=True,
                    )
                    import plotly.graph_objects as go
                    _ex = [
                        {"label": "Segment B — unreliable", "lift": 2.0, "ci": 3.5, "colour": "#F8696B"},
                        {"label": "Segment A — reliable", "lift": 5.0, "ci": 2.0, "colour": "#63BE7B"},
                    ]
                    _fig_ci = go.Figure()
                    for e in _ex:
                        lo, hi = e["lift"] - e["ci"], e["lift"] + e["ci"]
                        _fig_ci.add_trace(go.Scatter(
                            x=[lo, hi], y=[e["label"], e["label"]],
                            mode="lines", line=dict(width=8, color=e["colour"]),
                            showlegend=False, hoverinfo="skip",
                        ))
                        _fig_ci.add_trace(go.Scatter(
                            x=[e["lift"]], y=[e["label"]],
                            mode="markers+text", marker=dict(size=14, color="#fff", line=dict(width=2, color=e["colour"])),
                            text=[f"{e['lift']:.1f}%"], textposition="top center",
                            textfont=dict(size=13, color="#fafafa"),
                            showlegend=False, hoverinfo="skip",
                        ))
                        _fig_ci.add_annotation(
                            x=lo, y=e["label"], text=f"{lo:.1f}%", showarrow=False,
                            yshift=-18, font=dict(size=10, color="#bbb"),
                        )
                        _fig_ci.add_annotation(
                            x=hi, y=e["label"], text=f"{hi:.1f}%", showarrow=False,
                            yshift=-18, font=dict(size=10, color="#bbb"),
                        )
                    _fig_ci.add_vline(x=0, line_dash="dash", line_color="#888", line_width=1.5,
                                      annotation_text="zero", annotation_position="top",
                                      annotation_font_size=10, annotation_font_color="#999")
                    _fig_ci.update_layout(
                        height=200, margin=dict(l=10, r=10, t=30, b=10),
                        paper_bgcolor="rgba(0,0,0,0)", plot_bgcolor="rgba(0,0,0,0)",
                        xaxis=dict(title="Balance Lift %", zeroline=False, gridcolor="#333",
                                   ticksuffix="%", color="#ccc"),
                        yaxis=dict(color="#ccc"),
                        font=dict(color="#fafafa"),
                    )
                    st.plotly_chart(_fig_ci, width='stretch')

            # ── Recommended Audiences ────────────────────────────────────────────
            st.divider()
            st.markdown("### Recommended Audiences")
            st.caption(
                "Optimal segment groupings for each communication, ranked by lift. "
                "Use these as ready-made targeting lists — no manual segment selection needed."
            )
            with st.expander("How the audience is built", expanded=False):
                st.markdown("""
    **Step 1 — AND filter (optional)**  
    If you pick any *AND — mandatory segments*, only customers who belong to **all** of them are kept. Everything downstream (audience size, lift, N) is computed on this restricted pool.

    **Step 2 — Lift per segment**  
    Each segment's lift is the simple difference between the mean outcome of treated customers in that segment versus matched control customers in the same segment. Segments are ranked highest lift first.

    **Step 3 — Rank-order greedy selection until Min Audience is reached**  
    Segments are added in descending lift order until the total unique customer count reaches *Min audience*. The highest-lift segments always come first.

    **Step 4 — NOT suppression (bottom segments excluded)**  
    The lowest-lift segments from the remainder are flagged as NOT (avoid). Any customer in **at least one** NOT segment is removed from the final audience — even if they also qualify via a top segment.

    **Step 5 — Expected Lift shown**  
    The displayed lift is computed directly on the **final recommended cohort**: each customer counted once, compared to the matched control group. This is the lift you should actually observe if those customers are contacted.
    """)
            _tbl_ra = tbl

            _ra_c1, _ra_c2, _ra_c3, _ra_c4 = st.columns([2, 2, 2, 1])
            with _ra_c1:
                ra_metric = st.radio(
                    "Optimise for", ["Balance lift", "Accounts lift"],
                    horizontal=True, key="ra_metric",
                )
            with _ra_c2:
                ra_comm = st.selectbox(
                    "Communication", ordered_comms, key="ra_comm",
                )
            with _ra_c3:
                _ra_all_opts = sorted(tbl.index.astype(str).tolist())
                ra_and_segs = st.multiselect(
                    "AND — mandatory segments",
                    options=_ra_all_opts,
                    default=[s for s in st.session_state.get("ra_and_segs", []) if s in _ra_all_opts],
                    key="ra_and_segs",
                    help="These segments are always included in the recommended audience regardless of lift.",
                )
            with _ra_c4:
                ra_min_aud = st.number_input(
                    "Min audience", min_value=0, value=1000, step=500, key="ra_min_aud",
                    help="Add top-lift segments until the total reached audience is at least this size.",
                )

            ra_suffix = "_lift_bal" if ra_metric == "Balance lift" else "_lift_acct"
            ra_label  = "Balance Lift %" if ra_metric == "Balance lift" else "Accounts Lift %"

            def _pct_fmt_ra(v): return _fmt_pct(v, na_str="—")

            _ra_shown_segs: List[str] = []  # top / recommended segments
            _ra_bot_segs:   List[str] = []  # bottom / avoid segments

            _ra_col    = f"{ra_comm}{ra_suffix}"
            _ra_ci_col = f"{ra_comm}{'_lift_bal_ci' if ra_metric == 'Balance lift' else '_lift_acct_ci'}"
            _ra_n_col  = f"{ra_comm}_n"
            _ra_placeholder = st.empty()
            if _ra_col in _tbl_ra.columns:
                _ra_placeholder.info(f"⏳ Building recommended audience for **{ra_comm}**…")
                # Naïve per-segment lift directly from the precomputed table
                _naive_comm = _tbl_ra[_ra_col].dropna()
                _naive_comm.index = _naive_comm.index.astype(str)
                _ci_series = (
                    _tbl_ra[_ra_ci_col].reindex(_naive_comm.index).astype(float)
                    if _ra_ci_col in _tbl_ra.columns
                    else pd.Series(np.nan, index=_naive_comm.index)
                )
                _ci_series.index = _ci_series.index.astype(str)

                _sorted_comm = _naive_comm.sort_values(ascending=False)

                # Per-segment unique-customer counts for this communication,
                # filtered to only customers who also appear in ALL AND segments (if any set)
                _ra_comm_df_all = df[(df["contact_flag"] == 1) & (df["communication"] == ra_comm)]
                _and_idx_pre = [s for s in ra_and_segs if s in _sorted_comm.index]
                if _and_idx_pre:
                    _and_str_pre = set(str(s) for s in _and_idx_pre)
                    _cust_seg_sets = (
                        _ra_comm_df_all[_ra_comm_df_all["nsegment"].astype(str).isin(_and_str_pre)]
                        .groupby("alpha_key")["nsegment"]
                        .apply(lambda x: set(x.astype(str)))
                    )
                    _valid_and_custs = _cust_seg_sets[_cust_seg_sets.apply(lambda s: _and_str_pre.issubset(s))].index
                    _ra_comm_df_filt = _ra_comm_df_all[_ra_comm_df_all["alpha_key"].isin(_valid_and_custs)]
                else:
                    _ra_comm_df_filt = _ra_comm_df_all

                _seg_n_raw = (
                    _ra_comm_df_filt
                    .groupby("nsegment")["alpha_key"].nunique()
                )
                _seg_n_raw.index = _seg_n_raw.index.astype(str)

                # AND segments — hard constraints (must be present in the table)
                _and_idx = _and_idx_pre

                # Sort segments by lift descending (no tiers — simpler and faster)
                _sort_df = pd.DataFrame({'lift': _sorted_comm})
                _sort_df = _sort_df.sort_values('lift', ascending=False)

                # Single-pass greedy: pre-determine NOT segments from the absolute bottom
                # of the pool (lowest lift), build the NOT customer set upfront, then
                # accumulate only NOT-suppressed customers in one loop — no second pass needed.
                _non_and = _sort_df[~_sort_df.index.isin(_and_idx)]
                _and_str_set = set(str(s) for s in _and_idx)
                _seg_custs_map: dict = {}
                for _s in _and_idx:
                    _seg_custs_map[str(_s)] = set(
                        _ra_comm_df_filt[_ra_comm_df_filt["nsegment"].astype(str) == str(_s)]["alpha_key"]
                    )
                _remaining = list(_non_and.index)  # sorted by lift descending

                # Pre-load all customer sets up front
                for _seg in _remaining:
                    _s_str = str(_seg)
                    if _s_str not in _seg_custs_map:
                        _seg_custs_map[_s_str] = set(
                            _ra_comm_df_filt[_ra_comm_df_filt["nsegment"].astype(str) == _s_str]["alpha_key"]
                        )

                # Pre-determine NOT segments from the absolute bottom of the pool.
                # Building the NOT set upfront lets the greedy loop count NOT-suppressed
                # users at each step, so one pass always reaches the min-audience target.
                _bot_n = max(3, min(15, max(1, len(_remaining)) // 4))
                _ra_bot = _non_and.sort_values('lift').head(_bot_n)['lift']
                _not_segs_set = set(str(s) for s in _ra_bot.index)
                _not_custs: set = set()
                for _ns in _ra_bot.index:
                    _not_custs |= _seg_custs_map.get(str(_ns), set())

                # Greedy loop: skip NOT segments; accumulate only NOT-suppressed customers
                _running_custs: set = (
                    (set().union(*[_seg_custs_map[str(s)] for s in _and_idx]) - _not_custs)
                    if _and_idx else set()
                )
                _selected: List[str] = []
                for _seg in _remaining:
                    if str(_seg) in _not_segs_set:
                        continue
                    if len(_running_custs) >= ra_min_aud:
                        break
                    _selected.append(_seg)
                    _running_custs |= (_seg_custs_map[str(_seg)] - _not_custs)

                _ra_top_idx = list(dict.fromkeys(_and_idx + _selected))
                _ra_top     = _sorted_comm.reindex(
                    [s for s in _ra_top_idx if s in _sorted_comm.index], tolerance=None
                ).dropna()

                _final_custs = _running_custs  # already NOT-suppressed

                # Per-segment N (still used for the segment-level table display)
                _ra_n_vals = pd.Series({
                    s: len(_seg_custs_map.get(str(s), set()) - _not_custs)
                    for s in _ra_top.index
                }, dtype=float)
                _ra_n_vals.index = _ra_top.index
                _ra_users = len(_final_custs)

                # ── Cohort-level lift (direct, per-customer) ─────────────────────────
                # Each recommended customer is counted exactly ONCE regardless of how
                # many selected segments they belong to.  This is the only approach that
                # correctly handles multi-segment users and gives a lift estimate that
                # matches what you will observe when those exact customers are contacted.
                #
                # Method:
                #   treated cohort  = final recommended customers (deduplicated to 1 row/user)
                #   control cohort  = control customers in the same segment pool, same NOT
                #                     suppression applied (deduplicated to 1 row/user)
                #   lift            = treated_mean − control_mean  (two-sample difference)
                # ─────────────────────────────────────────────────────────────────────
                _suffix_col = "balance_pct_change" if ra_suffix == "_lift_bal" else "accounts_pct_change"

                # Treated side: filter to final audience, one row per user
                _treated_cohort = (
                    _ra_comm_df_filt[_ra_comm_df_filt["alpha_key"].isin(_final_custs)]
                    .drop_duplicates(subset="alpha_key")
                )

                # Control side: control customers — restricted to the same AND population
                # as the treated side, in selected segments, minus NOT-suppressed.
                # Critical: if AND segments are set, control customers must also be in all
                # AND segments, otherwise treated vs. control compare different populations.
                _ctrl_raw = df[(df["control_flag"] == 1) & (df["communication"] == ra_comm)]
                if _and_idx:
                    _and_str_set_ctrl = set(str(s) for s in _and_idx)
                    _ctrl_cust_segs = (
                        _ctrl_raw[_ctrl_raw["nsegment"].astype(str).isin(_and_str_set_ctrl)]
                        .groupby("alpha_key")["nsegment"]
                        .apply(lambda x: set(x.astype(str)))
                    )
                    _valid_and_ctrl = _ctrl_cust_segs[
                        _ctrl_cust_segs.apply(lambda s: _and_str_set_ctrl.issubset(s))
                    ].index
                    _ctrl_raw = _ctrl_raw[_ctrl_raw["alpha_key"].isin(_valid_and_ctrl)]
                _ctrl_sel_custs: set = set()
                for _s in _ra_top_idx:
                    _ctrl_sel_custs |= set(
                        _ctrl_raw[_ctrl_raw["nsegment"].astype(str) == str(_s)]["alpha_key"]
                    )
                _not_custs_ctrl: set = set()
                for _bot_s in _ra_bot.index:
                    _not_custs_ctrl |= set(
                        _ctrl_raw[_ctrl_raw["nsegment"].astype(str) == str(_bot_s)]["alpha_key"]
                    )
                _ctrl_cohort = (
                    _ctrl_raw[_ctrl_raw["alpha_key"].isin(_ctrl_sel_custs - _not_custs_ctrl)]
                    .drop_duplicates(subset="alpha_key")
                )

                _t_mean = _treated_cohort[_suffix_col].mean()
                _c_mean = _ctrl_cohort[_suffix_col].mean() if not _ctrl_cohort.empty else np.nan
                if pd.notna(_t_mean) and pd.notna(_c_mean):
                    _ra_w_lift = float(_t_mean - _c_mean)
                    _n_t, _n_c = len(_treated_cohort), len(_ctrl_cohort)
                    _se_t = float(_treated_cohort[_suffix_col].std(ddof=1)) / np.sqrt(_n_t) if _n_t >= 2 else np.nan
                    _se_c = float(_ctrl_cohort[_suffix_col].std(ddof=1))   / np.sqrt(_n_c) if _n_c >= 2 else np.nan
                    _ra_lift_ci = float(Z95 * np.sqrt(_se_t**2 + _se_c**2)) if (pd.notna(_se_t) and pd.notna(_se_c)) else np.nan
                elif pd.notna(_t_mean):
                    _ra_w_lift, _ra_lift_ci = float(_t_mean), np.nan
                else:
                    _ra_w_lift, _ra_lift_ci = np.nan, np.nan

                # Total customers reachable given the AND constraint (ceiling for greedy loop)
                _and_pool_total = int(_ra_comm_df_filt["alpha_key"].nunique())
                _rm1, _rm2 = st.columns(2)
                _rm1.metric("Recommended audience", f"{_ra_users:,} customers")
                _lift_display = f"{_ra_w_lift:.2%}" if pd.notna(_ra_w_lift) else "—"
                if pd.notna(_ra_lift_ci):
                    _lift_display += f"  ±{_ra_lift_ci:.2%}"
                _rm2.metric(
                    f"Expected {ra_label}",
                    _lift_display,
                    help="Treated audience mean minus matched-control mean — each customer counted once.",
                )
                if _and_idx and _and_pool_total < ra_min_aud:
                    st.warning(
                        f"⚠️ The AND constraint limits the total available audience to **{_and_pool_total:,} customers** "
                        f"for *{ra_comm}* — the min audience target of **{ra_min_aud:,}** cannot be reached. "
                        f"All {_and_pool_total:,} qualifying customers have been selected. "
                        f"Remove AND segments or lower the min audience to proceed normally.",
                        icon=None,
                    )
                _ci_label = f"{ra_label} ±CI"

                def _make_single_comm_table(series):
                    # Display the same direct (naïve) lift as the Data table above for consistency.
                    # OLS marginal lifts are used only for internal ranking/greedy selection.
                    _naive_disp = _naive_comm.reindex(series.index.astype(str))
                    _ci_v = (
                        _tbl_ra[_ra_ci_col].rename(index=str).reindex(series.index.astype(str)).values
                        if _ra_ci_col in _tbl_ra.columns
                        else [np.nan] * len(series)
                    )
                    _n_v  = [int(_tbl_ra.at[s, _ra_n_col]) if _ra_n_col in _tbl_ra.columns and s in _tbl_ra.index else 0 for s in series.index]
                    _d = pd.DataFrame({
                        "Description": [SEGMENT_DESCRIPTIONS.get(str(s), SEGMENT_LABELS.get(str(s), "")) for s in series.index],
                        ra_label:      _naive_disp.values,
                        _ci_label:     _ci_v,
                        "N":           _n_v,
                    }, index=series.index)
                    _d.index.name = None
                    return _d.style.format({ra_label: _pct_fmt_ra, _ci_label: _pct_fmt_ra, "N": "{:,.0f}"}, na_rep="")\
                        .apply(lambda s: s.map(_rdylgn), subset=[ra_label], axis=0)\
                        .apply(lambda s: s.map(_n_color), subset=["N"], axis=0)

                def _two_tables_html(items, height):
                    """Render multiple (label, styler) pairs into one HTML doc — zero spacing between tables."""
                    import re as _re3
                    _css = (
                        "* { box-sizing: border-box; }"
                        "body { font-family: -apple-system, BlinkMacSystemFont, sans-serif;"
                        "       font-size: 12px; background: transparent; color: #fafafa; margin: 0; padding: 4px; }"
                        "table { border-collapse: collapse; width: 100%; table-layout: fixed; }"
                        "thead th { background: #262730; color: #fafafa; padding: 6px 10px;"
                        "           text-align: left; position: sticky; top: 0; z-index: 2;"
                        "           border-bottom: 2px solid #555; white-space: nowrap; font-size: 11px; }"
                        "tbody td { padding: 4px 10px; border-bottom: 1px solid #333; white-space: nowrap; color: #111; }"
                        "tbody td:first-child { white-space: normal; word-break: break-word; }"
                        "tbody tr:hover td { outline: 1px solid #666; }"
                        "th.row_heading { background: #1c1e2a !important; font-size: 11px;"
                        "                 color: #ccc !important; font-weight: normal; cursor: help;"
                        "                 padding: 4px 6px 4px 6px; min-width: 50px; }"
                        "th.blank { background: #262730 !important; }"
                        ".tbl-lbl { font-size: 11px; font-weight: 600; color: #bbb;"
                        "           padding: 6px 2px 2px 2px; margin-top: 2px; }"
                        ".tbl-lbl:first-child { margin-top: 0; padding-top: 2px; }"
                    )
                    _blocks = []
                    for _item in items:
                        _lbl = _item[0]; _styler = _item[1]
                        _warn_p = _item[2] if len(_item) > 2 else None
                        _h = _styler.to_html()
                        if _warn_p:
                            _h = _inject_warn_flags(_h, _styler.data, _warn_p)
                        def _tip(m):
                            _tp = m.group(1); _cv = m.group(2); _id = _cv.strip()
                            _l = (SEGMENT_LABELS or {}).get(_id, ""); _d2 = (SEGMENT_DESCRIPTIONS or {}).get(_id, "")
                            _t = f"{_l} \u2014 {_d2}" if _l and _d2 else (_l or _d2)
                            if not _t: return m.group(0)
                            return f'{_tp} data-tip="{_t.replace(chr(34), chr(39))}">{_cv}</th>'
                        _h = _re3.sub(r'(<th[^>]*class="[^"]*row_heading[^"]*"[^>]*)>([^<]*)</th>', _tip, _h)
                        # Inject colgroup so every table in this iframe uses identical column widths
                        _ncols_inj = 1 + len(_styler.data.columns)  # index col + data cols
                        _pct_map = {
                            5: [10, 46, 16, 18, 10],
                            4: [14, 42, 25, 19],
                            6: [10, 30, 16, 18, 16, 10],
                        }
                        _col_pcts = _pct_map.get(_ncols_inj, [round(100 / _ncols_inj)] * _ncols_inj)
                        _cg = '<colgroup>' + ''.join(f'<col style="width:{p}%">' for p in _col_pcts) + '</colgroup>'
                        _h = _re3.sub(r'(<table[^>]*>)', r'\1' + _cg, _h, count=1)
                        _blocks.append(f'<div class="tbl-lbl">{_lbl}</div>{_h}')
                    _body = "".join(_blocks)
                    _js = (
                        "(function(){var tt=document.createElement('div');"
                        "tt.style.cssText='position:fixed;background:#1e2030;color:#eee;font-size:11px;"
                        "padding:5px 9px;border-radius:4px;border:1px solid #555;z-index:99999;"
                        "pointer-events:none;display:none;max-width:340px;word-wrap:break-word;"
                        "line-height:1.5;white-space:normal;box-shadow:0 2px 8px rgba(0,0,0,.5);';"
                        "document.body.appendChild(tt);"
                        "document.querySelectorAll('[data-tip]').forEach(function(el){"
                        "el.addEventListener('mouseenter',function(){tt.textContent=el.getAttribute('data-tip');tt.style.display='block';});"
                        "el.addEventListener('mousemove',function(e){tt.style.left=(e.clientX+14)+'px';tt.style.top=(e.clientY+14)+'px';});"
                        "el.addEventListener('mouseleave',function(){tt.style.display='none';});"
                        "})})()"
                    )
                    return (
                        f"<!DOCTYPE html><html><head><meta charset='utf-8'>"
                        f"<style>{_css}</style></head><body>"
                        f"<div style='overflow:auto;max-height:{height-20}px;'>{_body}</div>"
                        f"<script>{_js}</script></body></html>"
                    )

                _ra_styler2   = _make_single_comm_table(_ra_top)
                _ra_bot_styler = _make_single_comm_table(_ra_bot)
                _ra_shown_segs = [str(s) for s in _ra_top.index]
                _ra_bot_segs   = [str(s) for s in _ra_bot.index]

                _combined_h = max(300, min(900, 80 + (len(_ra_top) + len(_ra_bot)) * 30))
                components.html(
                    _two_tables_html([
                        (f"Top {len(_ra_top)} segments (highest lift)", _ra_styler2, [(ra_label, _ci_label)]),
                        (f"Bottom {len(_ra_bot)} segments (lowest lift — suppressed from final audience)", _ra_bot_styler, [(ra_label, _ci_label)]),
                    ], _combined_h),
                    height=_combined_h, scrolling=True,
                )
                _ra_placeholder.empty()
            else:
                _ra_placeholder.empty()
                st.info(f"No {ra_label} data available for **{ra_comm}**.")

            # ── Send recommended segments to Simulator ───────────────────────
            if _ra_shown_segs:
                st.markdown(
                    '<style>div:has(>#_ra_send_anchor)+div[data-testid="element-container"]'
                    '{margin-top:-28px!important}</style><div id="_ra_send_anchor"></div>',
                    unsafe_allow_html=True,
                )
                _, _ra_btn_col, _ = st.columns([1, 2, 1])
                if _ra_btn_col.button("Send to Audience Simulator", key="ra_send_sim", use_container_width=True):
                    st.session_state["sim_segs"]          = _ra_shown_segs
                    st.session_state["sim_segs_and"]      = list(_and_idx)
                    st.session_state["sim_segs_excl"]     = _ra_bot_segs
                    st.session_state["sim_run_triggered"] = True
                    _and_note = f", {len(_and_idx)} mandatory AND" if _and_idx else ""
                    st.success(f"Sent {len(_ra_shown_segs)} segment{'s' if len(_ra_shown_segs) != 1 else ''} to the Audience Simulator{_and_note} (+ {len(_ra_bot_segs)} excluded).")

        # ── Audience Demographics (all customers in audience_profile.csv) ─────────
        if _aud_df is not None and len(_aud_df) > 0:
            st.divider()
            st.subheader("Audience Demographics — All Customers")
            st.caption(f"{len(_aud_df):,} customers in audience_profile.csv")
            import plotly.graph_objects as go
            from scipy.stats import gaussian_kde as _gkde_dt

            def _kde_dt(series, nbins, line_color):
                _v = series.dropna().values
                if len(_v) < 5:
                    return None
                _, _edges = np.histogram(_v, bins=nbins)
                _bw = _edges[1] - _edges[0]
                _fn = _gkde_dt(_v)
                _xs = np.linspace(_v.min(), _v.max(), 300)
                _ys = _fn(_xs) * len(_v) * _bw
                return go.Scatter(x=_xs, y=_ys, mode="lines",
                                  line=dict(color=line_color, width=2.5),
                                  showlegend=False, name="")

            _d = _aud_df.copy()

            # Row 1: Age + Gender
            _d1, _d2 = st.columns([3, 2])
            with _d1:
                _age_bins2   = [18, 25, 35, 45, 55, 65, 75, 120]
                _age_labels2 = ["18-24", "25-34", "35-44", "45-54", "55-64", "65-74", "75+"]
                _d["_age_group2"] = pd.cut(_d["age"], bins=_age_bins2, labels=_age_labels2, right=False)
                _age_dist2 = (_d["_age_group2"].value_counts().reindex(_age_labels2).fillna(0).reset_index())
                _age_dist2.columns = ["Age group", "Customers"]
                _fig_age2 = px.bar(_age_dist2, x="Age group", y="Customers",
                                   title="Age Distribution", color="Customers",
                                   color_continuous_scale="Blues")
                _fig_age2.add_scatter(x=_age_dist2["Age group"], y=_age_dist2["Customers"],
                                      mode="lines+markers",
                                      line=dict(color="#103060", width=2, shape="spline", smoothing=1.0),
                                      marker=dict(size=6, color="#103060"), showlegend=False, name="")
                _age_mean2 = float(_d["age"].dropna().mean())
                _age_mean_bin2 = str(pd.cut([_age_mean2], bins=_age_bins2, labels=_age_labels2, right=False)[0])
                _fig_age2.add_shape(
                    type="line", xref="x", yref="paper",
                    x0=_age_mean_bin2, x1=_age_mean_bin2, y0=0, y1=1,
                    line=dict(color="red", width=2, dash="dash"),
                )
                _fig_age2.add_annotation(
                    x=_age_mean_bin2, yref="paper", y=1.05,
                    text=f"Mean: {_age_mean2:.1f}y", showarrow=False,
                    xanchor="left", font=dict(color="red", size=11),
                )
                _fig_age2.update_layout(height=300, margin=dict(l=20, r=20, t=50, b=30),
                                        coloraxis_showscale=False, showlegend=False)
                st.plotly_chart(_fig_age2, width='stretch')
            with _d2:
                _gender_counts2 = _d["gender"].fillna("Missing").value_counts().reset_index()
                _gender_counts2.columns = ["Gender", "Count"]
                _g_cmap2 = {"Male": "#89C4E1", "Female": "#FFB7C5", "Missing": "#CCCCCC"}
                _fig_gender2 = go.Figure(go.Pie(
                    labels=_gender_counts2["Gender"], values=_gender_counts2["Count"],
                    hole=0.4,
                    marker=dict(colors=[_g_cmap2.get(g, "#DDDDDD") for g in _gender_counts2["Gender"]]),
                    textinfo="percent+label"))
                _fig_gender2.update_layout(title=dict(text="Gender"), height=300,
                                           margin=dict(l=20, r=20, t=50, b=30))
                st.plotly_chart(_fig_gender2, width='stretch')

            st.markdown("<br>", unsafe_allow_html=True)

            # Row 2: Tenure + Deposit + Top 10 States
            _d3, _d4, _d5 = st.columns(3)
            with _d3:
                _fig_ten2 = px.histogram(_d, x="tenure_years", nbins=20, title="Tenure Distribution",
                                         labels={"tenure_years": "Tenure (years)"},
                                         color_discrete_sequence=["#4C9BE8"])
                _k2 = _kde_dt(_d["tenure_years"], 20, "#1a3a6b")
                if _k2:
                    _fig_ten2.add_trace(_k2)
                _ten_med2 = float(_d["tenure_years"].dropna().median())
                _fig_ten2.add_vline(x=_ten_med2, line=dict(color="red", width=2, dash="dash"),
                                    annotation_text=f"Median: {_ten_med2:.1f}y",
                                    annotation_position="top right",
                                    annotation_font=dict(color="red", size=11))
                _fig_ten2.update_layout(height=300, margin=dict(l=20, r=20, t=50, b=30))
                st.plotly_chart(_fig_ten2, width='stretch')
            with _d4:
                _dep_p99 = float(_d["amount_deposit_spot_balance"].dropna().quantile(0.99))
                _fig_dep2 = px.histogram(_d[_d["amount_deposit_spot_balance"] <= _dep_p99],
                                         x="amount_deposit_spot_balance", nbins=25, title="Deposit Balance",
                                         labels={"amount_deposit_spot_balance": "Balance ($)"},
                                         color_discrete_sequence=["#F4A261"])
                _k3 = _kde_dt(_d["amount_deposit_spot_balance"], 25, "#7a3100")
                if _k3:
                    _fig_dep2.add_trace(_k3)
                _dep_med2 = float(_d["amount_deposit_spot_balance"].dropna().median())
                _fig_dep2.add_vline(x=_dep_med2, line=dict(color="red", width=2, dash="dash"),
                                    annotation_text=f"Median: ${_dep_med2:,.0f}",
                                    annotation_position="top right",
                                    annotation_font=dict(color="red", size=11))
                _fig_dep2.update_xaxes(range=[0, _dep_p99])
                _fig_dep2.update_layout(height=300, margin=dict(l=20, r=20, t=50, b=30))
                st.plotly_chart(_fig_dep2, width='stretch')
            with _d5:
                _state_cnt2 = _d["state"].value_counts().head(10).reset_index()
                _state_cnt2.columns = ["State", "Customers"]
                _fig_state2 = px.bar(_state_cnt2, x="Customers", y="State", orientation="h",
                                     title="Top 10 States", color="Customers", color_continuous_scale="Teal")
                _fig_state2.update_layout(height=300, margin=dict(l=20, r=20, t=50, b=30),
                                          coloraxis_showscale=False,
                                          yaxis=dict(categoryorder="total ascending"))
                st.plotly_chart(_fig_state2, width='stretch')

            st.markdown("<br>", unsafe_allow_html=True)

            # Row 3: Product ownership rates + # products held
            _d6, _d7 = st.columns(2)
            with _d6:
                _prod_cols_dt = [c for c in PRODUCT_COLS if c in _d.columns]
                if _prod_cols_dt:
                    _prod_rates2 = ((_d[_prod_cols_dt].mean() * 100).sort_values(ascending=True).reset_index())
                    _prod_rates2.columns = ["Product", "% with product"]
                    _fig_prods2 = px.bar(_prod_rates2, x="% with product", y="Product", orientation="h",
                                         title="Product Ownership Rates", color="% with product",
                                         color_continuous_scale="Purpor",
                                         text=_prod_rates2["% with product"].map(lambda v: f"{v:.1f}%"))
                    _fig_prods2.update_traces(textposition="outside")
                    _fig_prods2.update_layout(height=340, margin=dict(l=20, r=80, t=40, b=20),
                                              coloraxis_showscale=False,
                                              xaxis=dict(ticksuffix="%", range=[0, 100]))
                    st.plotly_chart(_fig_prods2, width='stretch')
            with _d7:
                _nprod_cnt2 = _d["n_products"].value_counts().sort_index().reset_index()
                _nprod_cnt2.columns = ["Products", "Customers"]
                _fig_np2 = px.bar(_nprod_cnt2, x="Products", y="Customers", title="# Products Held",
                                   color="Customers", color_continuous_scale="YlOrRd")
                _fig_np2.add_scatter(x=_nprod_cnt2["Products"], y=_nprod_cnt2["Customers"],
                                      mode="lines+markers",
                                      line=dict(color="#5a0000", width=2, shape="spline", smoothing=0.8),
                                      marker=dict(size=6, color="#5a0000"), showlegend=False, name="")
                _nprod_mean2 = float(_d["n_products"].dropna().mean())
                _fig_np2.add_vline(x=_nprod_mean2, line=dict(color="red", width=2, dash="dash"),
                                   annotation_text=f"Mean: {_nprod_mean2:.1f}",
                                   annotation_position="top right",
                                   annotation_font=dict(color="red", size=11))
                _fig_np2.update_layout(height=340, margin=dict(l=20, r=20, t=50, b=30),
                                        coloraxis_showscale=False)
                st.plotly_chart(_fig_np2, width='stretch')

            st.markdown("<br>", unsafe_allow_html=True)

            # Row 4: Deposits vs IXI scatter (full width, capped at 99th pct)
            _sow_df2 = _d[["amount_deposit_spot_balance", "total_deposits_ixi", "sow"]].dropna() if "total_deposits_ixi" in _d.columns else None
            if _sow_df2 is not None and len(_sow_df2) > 0:
                _sow_df2 = _sow_df2[
                    (_sow_df2["amount_deposit_spot_balance"] <= 200_000) &
                    (_sow_df2["total_deposits_ixi"] <= 1_500_000)
                ]
                _sow_overall_mean2 = float(_sow_df2["sow"].mean())
                _fig_sow2 = px.scatter(
                    _sow_df2, x="amount_deposit_spot_balance", y="total_deposits_ixi",
                    color="sow", color_continuous_scale="RdYlGn", range_color=[0, 1],
                    title=f"Deposits VS IXI (SoW {_sow_overall_mean2:.1%})",
                    labels={"amount_deposit_spot_balance": "Deposit ($)",
                            "total_deposits_ixi": "IXI ($)", "sow": "SoW"},
                    opacity=0.5)
                _fig_sow2.update_xaxes(range=[0, 200_000])
                _fig_sow2.update_yaxes(range=[0, 1_500_000])
                # ── Percentile contour lines (add_shape layer='above' → always on top) ──
                if len(_sow_df2) >= 20:
                    import plotly.graph_objects as go
                    import matplotlib
                    matplotlib.use('Agg')
                    import matplotlib.pyplot as _plt2
                    from scipy.stats import gaussian_kde as _gkde_sow2
                    from scipy.ndimage import gaussian_filter as _gf2
                    _x2 = _sow_df2["amount_deposit_spot_balance"].values
                    _y2 = _sow_df2["total_deposits_ixi"].values
                    try:
                        _kde2 = _gkde_sow2(np.vstack([_x2, _y2]), bw_method=0.3)
                        _xi2 = np.linspace(_x2.min(), _x2.max(), 200)
                        _yi2 = np.linspace(_y2.min(), _y2.max(), 200)
                        _XX2, _YY2 = np.meshgrid(_xi2, _yi2)
                        _ZZ2 = _kde2(np.vstack([_XX2.ravel(), _YY2.ravel()])).reshape(_XX2.shape)
                        _ZZ2 = _gf2(_ZZ2, sigma=3)
                        _z_pts2 = _kde2(np.vstack([_x2, _y2]))
                        _sow_vals2 = _sow_df2["sow"].values
                        for _plbl2, _clr2 in [
                            (90, "#2166ac"),   # blue  — outermost (90% of pts inside)
                            (50, "#1a9850"),   # green — 50% inside
                            (25, "#f4a11d"),   # amber — 25% inside
                            (10, "#c0392b"),   # red   — dense core
                        ]:
                            _lvl2 = float(np.percentile(_z_pts2, 100 - _plbl2))
                            if _lvl2 <= 0:
                                continue
                            _mask2 = _z_pts2 >= _lvl2
                            _sow_grp2 = float(np.mean(_sow_vals2[_mask2])) if _mask2.any() else float('nan')
                            _sow_lbl2 = f" (SoW {_sow_grp2:.1%})" if not np.isnan(_sow_grp2) else ""
                            # Dummy trace for legend entry only (shapes don't appear in legend)
                            _fig_sow2.add_trace(go.Scatter(
                                x=[None], y=[None], mode='lines',
                                line=dict(color=_clr2, width=2.5),
                                name=f'{_plbl2}th pct{_sow_lbl2}',
                                legendgroup=f'pct{_plbl2}_2', showlegend=True,
                            ))
                            _mfig2, _max2 = _plt2.subplots()
                            _cs2 = _max2.contour(_XX2, _YY2, _ZZ2, levels=[_lvl2])
                            _plt2.close(_mfig2)
                            for _seg2 in _cs2.get_paths():
                                _verts2 = _seg2.vertices
                                if len(_verts2) < 5:
                                    continue
                                # Open contour at lower-left corner, extend to both axes
                                _myi2 = int(np.argmin(_verts2[:, 1]))
                                _mxi2 = int(np.argmin(_verts2[:, 0]))
                                _n2 = len(_verts2)
                                _rol2 = np.roll(_verts2, -_myi2, axis=0)
                                _nmx2 = (_mxi2 - _myi2) % _n2
                                _arc_a2 = _rol2[:_nmx2 + 1]
                                _arc_b2 = _rol2[_nmx2:]
                                if _arc_a2.shape[0] < 3 or _arc_b2.shape[0] < 3:
                                    _opn2 = _rol2
                                elif np.mean(_arc_a2[:, 0] + _arc_a2[:, 1]) >= np.mean(_arc_b2[:, 0] + _arc_b2[:, 1]):
                                    _opn2 = _arc_a2
                                else:
                                    _opn2 = _arc_b2
                                _pts2 = np.vstack([
                                    [max(float(_opn2[0, 0]), 0), 0],
                                    _opn2,
                                    [0, max(float(_opn2[-1, 1]), 0)],
                                ])
                                _svg2 = "M " + " L ".join(f"{float(_px):.4f},{float(_py):.4f}" for _px, _py in _pts2)
                                _fig_sow2.add_shape(
                                    type='path', path=_svg2,
                                    xref='x', yref='y',
                                    line=dict(color=_clr2, width=2.5),
                                    fillcolor='rgba(0,0,0,0)',
                                    layer='above',
                                )
                    except Exception:
                        pass
                _fig_sow2.update_layout(height=380, margin=dict(l=20, r=20, t=50, b=30),
                                        legend=dict(x=0.01, y=0.99, bgcolor='rgba(255,255,255,0.8)',
                                                    bordercolor='#aaa', borderwidth=1))
                st.plotly_chart(_fig_sow2, width='stretch')


    # ══════════════════════════════════════════════════════════
    # AUDIENCE SIMULATOR TAB
    # ══════════════════════════════════════════════════════════
    with tab_simulator:
        st.subheader("Audience Performance Simulator")
        st.caption(
            "Select any set of segments below and we'll estimate the **expected balance lift per "
            "communication** if you sent to only those customers. "
            "Each customer is counted once — the lift shown is what you'd actually observe on that cohort."
        )

        with st.expander("How the numbers are calculated", expanded=False):
            st.markdown("""
    **Expected lift (headline KPI)**  
    Each customer in your selected segments is counted **exactly once**, regardless of how many of those segments they belong to.  
    Lift = mean outcome for treated customers in this cohort minus mean outcome for matched control customers in the same segment pool.  
    This is the same methodology as the Recommendation tab — the number you see is what you would actually have observed historically if you had contacted exactly those people.

    **Per-segment breakdown table**  
    The table below shows per-segment figures (naïve group means) for reference. These are not used in the headline KPI — they are informational only.

    **Projected absolute € increase**  
    Formula: `unique customers × mean incremental balance change (treated − control baseline)`  
    The incremental change is the direct observed average € change per customer in the treated cohort minus the control cohort baseline — no pct × avg-balance reconstruction.  
    Note: this is an *expected* estimate based on historical data — actual results will vary.

    **Projected account openings**  
    Same logic: `unique customers × mean incremental accounts change (treated − control baseline)`.
            """)

        if "tbl" not in dir() or tbl.empty:
            st.warning("No table data — adjust filters in the Data tab first.")
        else:
            all_segs_sim = sorted(tbl.index.astype(str).tolist())

            sim_segs = st.multiselect(
                "OR — segments to include (union)",
                options=all_segs_sim,
                default=st.session_state.get("sim_segs", []),
                format_func=lambda x: x,
                key="sim_segs",
                max_selections=len(all_segs_sim),
                help="A customer qualifies if they belong to ANY of these segments.",
            )

            # AND multiselect — intersection filter
            sim_segs_and = st.multiselect(
                "AND — must also be in (intersection)",
                options=all_segs_sim,
                default=[s for s in st.session_state.get("sim_segs_and", []) if s in all_segs_sim],
                format_func=lambda x: x,
                key="sim_segs_and",
                max_selections=len(all_segs_sim),
                help="When non-empty, only segments that appear in BOTH the OR list and this list are kept. "
                     "Use this to narrow: e.g. OR=[A,B,C] AND=[B,C,D] → effective segments are {B,C}.",
            )

            # NOT multiselect — exclusion
            sim_segs_excl = st.multiselect(
                "NOT — segments to exclude",
                options=all_segs_sim,
                default=[s for s in st.session_state.get("sim_segs_excl", []) if s in all_segs_sim],
                format_func=lambda x: x,
                key="sim_segs_excl",
                max_selections=len(all_segs_sim),
                help="Any segment listed here is removed from the final set, even if it appears above.",
            )
            # Effective segments: OR ∩ AND (if set) − NOT
            _sim_segs_eff = [
                s for s in (sim_segs or [])
                if (not sim_segs_and or s in sim_segs_and)
                and s not in sim_segs_excl
            ]

            if st.button("▶ Run simulation", key="sim_run_btn", type="primary"):
                st.session_state["sim_run_triggered"] = True
                st.session_state["_sim_segs_snapshot"] = list(_sim_segs_eff)
                st.session_state["_sim_excl_snapshot"] = list(sim_segs_excl or [])
            sim_metric = st.session_state.get("sim_metric", "Balance % lift")

            # Use the snapshot captured at button-click time for all computation.
            # This means editing the multiselects after a run does NOT re-trigger the
            # expensive loop — only clicking ▶ Run simulation does.
            _compute_segs = st.session_state.get("_sim_segs_snapshot", [])
            _compute_excl = st.session_state.get("_sim_excl_snapshot", [])

            if st.session_state.get("sim_run_triggered") and _compute_segs:
                # ── Pre-filter table rows and compute projected metrics ───────────
                sub_all = tbl.loc[tbl.index.isin(_compute_segs)].copy()
                # Total unique customers header: across all comms, with NOT suppression
                # scoped to treated customers (contact_flag=1) only.
                _global_not_ids: set = set(
                    df[(df["contact_flag"] == 1) & (df["nsegment"].isin(_compute_excl))]["alpha_key"]
                ) if _compute_excl else set()
                _sim_total_users_base = df[
                    (df["contact_flag"] == 1) & (df["nsegment"].isin(_compute_segs))
                ]["alpha_key"]
                if _global_not_ids:
                    _sim_total_users_base = _sim_total_users_base[~_sim_total_users_base.isin(_global_not_ids)]
                _sim_total_users = _sim_total_users_base.nunique()

                def _w_avg(col, n_vals):
                    if col not in sub_all.columns:
                        return np.nan
                    valid = sub_all[col].notna()
                    tot = n_vals[valid].sum()
                    return float((sub_all.loc[valid, col] * n_vals[valid]).sum() / tot) if tot > 0 else np.nan

                sim_summary_rows = []
                for comm in ordered_comms:
                    n_col = f"{comm}_n"
                    if n_col not in tbl.columns:
                        continue
                    n_vals  = sub_all[n_col].fillna(0) if n_col in sub_all.columns else pd.Series(dtype=float)
                    total_n = int(n_vals.sum())

                    # Cohort-level direct lift: each customer counted once regardless of
                    # how many selected segments they belong to (same method as RA tab).
                    # NOT customers are scoped to this specific communication so the count
                    # matches what Recommended Audiences shows for the same comm.
                    _comm_not_ids: set = set(
                        df[(df["contact_flag"] == 1) & (df["communication"] == comm)
                           & (df["nsegment"].isin(_compute_excl))]["alpha_key"]
                    ) if _compute_excl else set()
                    _comm_df_t = (
                        df[(df["contact_flag"] == 1) & (df["communication"] == comm)
                           & (df["nsegment"].isin(_compute_segs))]
                        .drop_duplicates("alpha_key")
                    )
                    if _comm_not_ids:
                        _comm_df_t = _comm_df_t[~_comm_df_t["alpha_key"].isin(_comm_not_ids)]
                    _comm_df_c = (
                        df[(df["control_flag"] == 1) & (df["communication"] == comm)
                           & (df["nsegment"].isin(_compute_segs))]
                        .drop_duplicates("alpha_key")
                    )
                    if _comm_not_ids:
                        _comm_df_c = _comm_df_c[~_comm_df_c["alpha_key"].isin(_comm_not_ids)]
                    _t_bal  = _comm_df_t["balance_pct_change"].mean()  if not _comm_df_t.empty else np.nan
                    _c_bal  = _comm_df_c["balance_pct_change"].mean()  if not _comm_df_c.empty else np.nan
                    _t_acct = _comm_df_t["accounts_pct_change"].mean() if not _comm_df_t.empty else np.nan
                    _c_acct = _comm_df_c["accounts_pct_change"].mean() if not _comm_df_c.empty else np.nan

                    w_lift_bal  = (float(_t_bal  - _c_bal)  if pd.notna(_t_bal)  and pd.notna(_c_bal)
                                   else _w_avg(f"{comm}_lift_bal",  n_vals))
                    w_lift_acct = (float(_t_acct - _c_acct) if pd.notna(_t_acct) and pd.notna(_c_acct)
                                   else _w_avg(f"{comm}_lift_acct", n_vals))
                    w_raw_bal   = float(_t_bal)  if pd.notna(_t_bal)  else _w_avg(f"{comm}_bal",  n_vals)
                    w_raw_acct  = float(_t_acct) if pd.notna(_t_acct) else _w_avg(f"{comm}_acct", n_vals)
                    w_lift = w_lift_bal  if sim_metric == "Balance % lift" else w_lift_acct
                    w_raw  = w_raw_bal   if sim_metric == "Balance % lift" else w_raw_acct

                    _comm_n_u = len(_comm_df_t)
                    # Projection: direct incremental absolute change (treated − control baseline)
                    # This avoids the inflation risk of pct×avg_balance when balance is skewed.
                    _avg_abs_bal_t  = _comm_df_t["balance_abs_change"].mean()  if not _comm_df_t.empty else np.nan
                    _avg_abs_bal_c  = _comm_df_c["balance_abs_change"].mean()  if not _comm_df_c.empty else np.nan
                    _avg_abs_acct_t = _comm_df_t["accounts_abs_change"].mean() if not _comm_df_t.empty else np.nan
                    _avg_abs_acct_c = _comm_df_c["accounts_abs_change"].mean() if not _comm_df_c.empty else np.nan
                    _incr_bal  = float(_avg_abs_bal_t  - _avg_abs_bal_c)  if (pd.notna(_avg_abs_bal_t)  and pd.notna(_avg_abs_bal_c))  else (float(_avg_abs_bal_t)  if pd.notna(_avg_abs_bal_t)  else np.nan)
                    _incr_acct = float(_avg_abs_acct_t - _avg_abs_acct_c) if (pd.notna(_avg_abs_acct_t) and pd.notna(_avg_abs_acct_c)) else (float(_avg_abs_acct_t) if pd.notna(_avg_abs_acct_t) else np.nan)
                    proj_eur   = _comm_n_u * _incr_bal  if (pd.notna(_incr_bal)  and _comm_n_u > 0) else np.nan
                    proj_accts = _comm_n_u * _incr_acct if (pd.notna(_incr_acct) and _comm_n_u > 0) else np.nan

                    sim_summary_rows.append({
                        "Communication":   comm,
                        "Expected Lift":   w_lift,
                        "Treatment Avg":   w_raw,
                        "Total N":         total_n,
                        "Comm Users":      _comm_n_u,
                        "Projected Bal \u20ac": proj_eur,
                        "Proj. Accounts":  proj_accts,
                        "_rank":           _rank(comm),
                    })

                sim_summary = (pd.DataFrame(sim_summary_rows)
                               .sort_values("_rank")
                               .drop(columns="_rank")
                               .reset_index(drop=True))

                _excl_note = f" ({len(_compute_excl)} excluded)" if _compute_excl else ""
                _tot_c1, _tot_c2 = st.columns([2, 1])
                _tot_c1.metric(
                    "Total unique customers (all selected segments)" + _excl_note,
                    f"{_sim_total_users:,}",
                    help="Unique people (alpha_key) who were contacted and appear in at least one of the selected segments.",
                )
                with _tot_c2:
                    st.radio("Metric", ["Balance % lift", "Accounts % lift"], key="sim_metric", horizontal=True)
                _lift_suffix = "_lift_bal" if sim_metric == "Balance % lift" else "_lift_acct"
                _raw_suffix  = "_bal"      if sim_metric == "Balance % lift" else "_acct"
                _y_label     = "Expected Balance Lift %" if sim_metric == "Balance % lift" else "Expected Accounts Lift %"

                # Show kpi metrics row
                st.markdown("#### Expected lift across selected segments")
                kpi_cols = st.columns(len(sim_summary))
                for i, row in sim_summary.iterrows():
                    lift_val = row["Expected Lift"]
                    label    = row["Communication"]
                    u_val    = int(row["Comm Users"]) if pd.notna(row.get("Comm Users")) else 0
                    delta_str = f"{u_val:,} customers"
                    _kpi_help = (
                        f"**{u_val:,} unique customers** received **{label}** and belong to at "
                        f"least one of the {len(_compute_segs)} selected segments. "
                        f"Each person is counted once. The projected figures below use this same headcount."
                    )
                    if pd.notna(lift_val):
                        kpi_cols[i].metric(label, f"{lift_val:.2%}", delta_str, help=_kpi_help)
                    else:
                        kpi_cols[i].metric(label, "—", delta_str, help=_kpi_help)

                # Projected row — only show the metric matching the selected metric
                if sim_metric == "Balance % lift":
                    st.markdown("#### Projected absolute balance increase")
                    proj_b_cols = st.columns(len(sim_summary))
                    for i, row in sim_summary.iterrows():
                        pv = row["Projected Bal \u20ac"]
                        proj_b_cols[i].metric(
                            row["Communication"],
                            f"\u20ac{pv:,.0f}" if pd.notna(pv) else "—",
                            help="Formula: unique customers × mean incremental balance change (treated − control baseline). Direct € observation, not pct × avg balance.",
                        )
                else:
                    st.markdown("#### Projected account openings")
                    proj_a_cols = st.columns(len(sim_summary))
                    for i, row in sim_summary.iterrows():
                        av = row["Proj. Accounts"]
                        proj_a_cols[i].metric(
                            row["Communication"],
                            f"{av:,.0f}" if pd.notna(av) else "—",
                            help="Formula: unique customers × mean incremental accounts change (treated − control baseline). Direct observation, not pct × avg accounts.",
                        )

                st.divider()

                # ── Per-segment detail table ──────────────────────────────────────
                _h_col, _t_col = st.columns([3, 1])
                _h_col.markdown("#### Per-segment breakdown")
                show_n_pct = _t_col.checkbox("Replace N with proportion", value=False, key="show_n_pct_sim")

                # Build columns: lift and N per comm (no CI)
                _intl_cols = []
                for c in ordered_comms:
                    if f"{c}{_lift_suffix}" in tbl.columns:
                        _intl_cols.append(f"{c}{_lift_suffix}")
                    if f"{c}_n" in tbl.columns:
                        _intl_cols.append(f"{c}_n")

                # Include all segments: OR (sim_segs), AND (sim_segs_and), NOT (sim_segs_excl)
                _all_detail_segs = list(dict.fromkeys(
                    list(sim_segs or []) + list(sim_segs_and or []) + list(sim_segs_excl or [])
                ))
                detail = tbl.loc[tbl.index.isin(_all_detail_segs), [c for c in _intl_cols if c in tbl.columns]].copy()
                detail.index.name = None

                _metric_lbl = " Bal%" if sim_metric == "Balance % lift" else " Acct%"
                col_rename_sim = {f"{c}{_lift_suffix}": f"{c}{_metric_lbl}" for c in ordered_comms}
                col_rename_sim.update({f"{c}_n": f"{c} N" for c in ordered_comms})
                detail = detail.rename(columns=col_rename_sim)

                # Role column
                _sim_segs_and_set  = set(str(s) for s in (sim_segs_and  or []))
                _sim_segs_excl_set = set(str(s) for s in (sim_segs_excl or []))
                def _get_role(seg):
                    s = str(seg)
                    if s in _sim_segs_excl_set: return "NOT"
                    if s in _sim_segs_and_set:  return "AND"
                    return "OR"
                detail.insert(0, "Role", [_get_role(s) for s in detail.index])

                # Sort: AND → OR → NOT
                _role_order = {"AND": 0, "OR": 1, "NOT": 2}
                detail = detail.assign(_rs=detail["Role"].map(_role_order)).sort_values("_rs").drop(columns="_rs")

                # N as proportion
                n_disp_cols = [f"{c} N" for c in ordered_comms if f"{c} N" in detail.columns]
                if show_n_pct:
                    for _nc in n_disp_cols:
                        _ncsum = detail[_nc].sum()
                        if _ncsum > 0:
                            detail[_nc] = detail[_nc] / _ncsum
                    n_fmt = {_nc: (lambda v: f"{v:.1%}" if pd.notna(v) else "") for _nc in n_disp_cols}
                else:
                    n_fmt = {_nc: "{:,.0f}" for _nc in n_disp_cols}

                lift_disp_cols = [f"{c}{_metric_lbl}" for c in ordered_comms if f"{c}{_metric_lbl}" in detail.columns]
                pct_fmt = {c: _fmt_pct for c in lift_disp_cols}
                fmt_all = {**pct_fmt, **n_fmt}

                def _lift_color(col_series):
                    return col_series.map(_rdylgn)

                def _n_color_sim(col_series):
                    return col_series.map(_n_color)

                def _role_color(v):
                    if v == "NOT": return "background-color: #ffdddd; color: #880000"
                    if v == "AND": return "background-color: #e8d5ff; color: #5500aa"
                    if v == "OR":  return "background-color: #d0e8ff; color: #003399"
                    return ""
                styler_sim = detail.style.format(fmt_all, na_rep="")
                if lift_disp_cols:
                    styler_sim = styler_sim.apply(_lift_color, subset=lift_disp_cols, axis=0)
                if n_disp_cols:
                    styler_sim = styler_sim.apply(_n_color_sim, subset=n_disp_cols, axis=0)
                styler_sim = styler_sim.apply(lambda s: s.map(_role_color), subset=["Role"], axis=0)
                _sim_h = max(300, min(600, 60 + len(detail) * 30))
                components.html(
                    _styled_html_table(styler_sim, SEGMENT_LABELS, SEGMENT_DESCRIPTIONS, height=_sim_h),
                    height=_sim_h, scrolling=True,
                )

                # ── Excel export — full breakdown table ───────────────────────────
                _dl_detail = detail.copy()
                if show_n_pct:
                    # Re-export with absolute N (not proportion) for clarity
                    for c in ordered_comms:
                        _nc = f"{c} N"
                        if _nc in _dl_detail.columns:
                            _ncsum_raw = tbl.loc[tbl.index.isin(_all_detail_segs), f"{c}_n"].sum() if f"{c}_n" in tbl.columns else 0
                            if _ncsum_raw > 0:
                                _dl_detail[_nc] = _dl_detail[_nc] * _ncsum_raw
                _dl_detail.insert(0, "Description", [SEGMENT_DESCRIPTIONS.get(str(s), "") for s in _dl_detail.index])
                _dl_detail.insert(0, "Label", [SEGMENT_LABELS.get(str(s), "") for s in _dl_detail.index])
                _dl_detail.index.name = "nsegment"
                st.markdown(
                    '<style>div:has(>#_dl_anchor)+div[data-testid="element-container"]'
                    '{margin-top:-16px!important}</style><div id="_dl_anchor"></div>',
                    unsafe_allow_html=True,
                )
                _, _dl_col, _ = st.columns([2, 2, 2])
                if _dl_col.button("Download Excel (.xlsx)", use_container_width=True,
                                   help="Opens a column selector before downloading."):
                    _sim_excel_dialog(_dl_detail, sim_summary, ordered_comms, _metric_lbl)
                # ── Audience Profile ──────────────────────────────────────────────────────
                st.markdown("---")
                st.markdown("#### Audience Profile")
                if _aud_df is None:
                    st.info(
                        "No audience profile data found. "
                        "Please provide an `audience_profile.csv` file."
                    )
                else:
                    # Resolve final user set: users in the selected (OR+AND) segments minus NOT
                    _aud_or_keys = df.loc[df["nsegment"].isin(_compute_segs), "alpha_key"]
                    if _compute_excl:
                        _aud_excl_keys = set(df.loc[df["nsegment"].isin(_compute_excl), "alpha_key"])
                        _aud_keys = _aud_or_keys[~_aud_or_keys.isin(_aud_excl_keys)].unique()
                    else:
                        _aud_keys = _aud_or_keys.unique()
                    _aud = _aud_df[_aud_df["alpha_key"].isin(_aud_keys)].copy()
                    if len(_aud) == 0:
                        st.info("No demographic records found for the selected segments.")
                    else:
                        # ── KPI row ──────────────────────────────────────────────
                        _pk1, _pk2, _pk3, _pk4, _pk5 = st.columns(5)
                        _pk1.metric("Audience size", f"{len(_aud):,}")
                        _pk2.metric("Median age", f"{_aud['age'].median():.0f} yrs")
                        _pk3.metric("Median tenure", f"{_aud['tenure_years'].median():.1f} yrs")
                        _pk4.metric(
                            "Mobile / ROB active",
                            f"{_aud['flag_last90_active_mob_rob'].mean() * 100:.0f}%",
                        )
                        _pk5.metric("Median SoW", f"{_aud['sow'].median() * 100:.0f}%")

                        import plotly.graph_objects as go
                        from scipy.stats import gaussian_kde as _gkde

                        def _kde_trace(series, nbins, bar_color, line_color):
                            """Return a KDE Scatter trace scaled to match a histogram's bar heights."""
                            _v = series.dropna().values
                            if len(_v) < 5:
                                return None
                            _, _edges = np.histogram(_v, bins=nbins)
                            _bw = _edges[1] - _edges[0]
                            _fn = _gkde(_v)
                            _xs = np.linspace(_v.min(), _v.max(), 300)
                            _ys = _fn(_xs) * len(_v) * _bw
                            return go.Scatter(
                                x=_xs, y=_ys, mode="lines",
                                line=dict(color=line_color, width=2.5),
                                showlegend=False, name="",
                            )

                        st.markdown("<br>", unsafe_allow_html=True)

                        # ── Row 1: Age distribution + Gender ─────────────────────
                        _ac1, _ac2 = st.columns([3, 2])
                        with _ac1:
                            _age_bins   = [18, 25, 35, 45, 55, 65, 75, 120]
                            _age_labels = ["18-24", "25-34", "35-44", "45-54", "55-64", "65-74", "75+"]
                            _aud["_age_group"] = pd.cut(
                                _aud["age"], bins=_age_bins, labels=_age_labels, right=False
                            )
                            _age_dist = (
                                _aud["_age_group"].value_counts()
                                .reindex(_age_labels)
                                .fillna(0)
                                .reset_index()
                            )
                            _age_dist.columns = ["Age group", "Customers"]
                            _fig_age = px.bar(
                                _age_dist, x="Age group", y="Customers",
                                title="Age Distribution",
                                color="Customers", color_continuous_scale="Blues",
                            )
                            _fig_age.add_scatter(
                                x=_age_dist["Age group"], y=_age_dist["Customers"],
                                mode="lines+markers",
                                line=dict(color="#103060", width=2, shape="spline", smoothing=1.0),
                                marker=dict(size=6, color="#103060"),
                                showlegend=False, name="",
                            )
                            _fig_age.update_layout(
                                height=300, margin=dict(l=20, r=20, t=50, b=30),
                                coloraxis_showscale=False, showlegend=False,
                            )
                            st.plotly_chart(_fig_age, width='stretch')
                        with _ac2:
                            _gender_vals = _aud["gender"].fillna("Missing")
                            _gender_counts = _gender_vals.value_counts().reset_index()
                            _gender_counts.columns = ["Gender", "Count"]
                            _g_cmap = {"Male": "#89C4E1", "Female": "#FFB7C5", "Missing": "#CCCCCC"}
                            _g_colors = [_g_cmap.get(g, "#DDDDDD") for g in _gender_counts["Gender"]]
                            _fig_gender = go.Figure(go.Pie(
                                labels=_gender_counts["Gender"],
                                values=_gender_counts["Count"],
                                hole=0.4,
                                marker=dict(colors=_g_colors),
                                textinfo="percent+label",
                            ))
                            _fig_gender.update_layout(
                                title=dict(text="Gender"),
                                height=300, margin=dict(l=20, r=20, t=50, b=30),
                                legend=dict(orientation="v", yanchor="middle", y=0.5),
                                showlegend=True,
                            )
                            st.plotly_chart(_fig_gender, width='stretch')

                        st.markdown("<br>", unsafe_allow_html=True)

                        # ── Row 2: Tenure + Deposit Balance + Top 10 States ──────
                        _ac3, _ac4, _ac5 = st.columns(3)
                        with _ac3:
                            _fig_ten = px.histogram(
                                _aud, x="tenure_years", nbins=20,
                                title="Tenure Distribution",
                                labels={"tenure_years": "Tenure (years)"},
                                color_discrete_sequence=["#4C9BE8"],
                            )
                            _kde_ten = _kde_trace(_aud["tenure_years"], 20, "#4C9BE8", "#1a3a6b")
                            if _kde_ten:
                                _fig_ten.add_trace(_kde_ten)
                            _fig_ten.update_layout(height=300, margin=dict(l=20, r=20, t=50, b=30))
                            st.plotly_chart(_fig_ten, width='stretch')
                        with _ac4:
                            _fig_dep = px.histogram(
                                _aud, x="amount_deposit_spot_balance", nbins=25,
                                title="Deposit Balance",
                                labels={"amount_deposit_spot_balance": "Balance ($)"},
                                color_discrete_sequence=["#F4A261"],
                            )
                            _kde_dep = _kde_trace(_aud["amount_deposit_spot_balance"], 25, "#F4A261", "#7a3100")
                            if _kde_dep:
                                _fig_dep.add_trace(_kde_dep)
                            _fig_dep.update_layout(height=300, margin=dict(l=20, r=20, t=50, b=30))
                            st.plotly_chart(_fig_dep, width='stretch')
                        with _ac5:
                            _state_cnt = _aud["state"].value_counts().head(10).reset_index()
                            _state_cnt.columns = ["State", "Customers"]
                            _fig_state = px.bar(
                                _state_cnt, x="Customers", y="State",
                                orientation="h", title="Top 10 States",
                                color="Customers", color_continuous_scale="Teal",
                            )
                            _fig_state.update_layout(
                                height=300, margin=dict(l=20, r=20, t=50, b=30),
                                coloraxis_showscale=False,
                                yaxis=dict(categoryorder="total ascending"),
                            )
                            st.plotly_chart(_fig_state, width='stretch')

                        st.markdown("<br>", unsafe_allow_html=True)

                        # ── Row 3: Product ownership rates + # products held ──────
                        _ac6, _ac7 = st.columns(2)
                        with _ac6:
                            _prod_cols_present = [c for c in PRODUCT_COLS if c in _aud.columns]
                            if _prod_cols_present:
                                _prod_rates = (
                                    (_aud[_prod_cols_present].mean() * 100)
                                    .sort_values(ascending=True)
                                    .reset_index()
                                )
                                _prod_rates.columns = ["Product", "% with product"]
                                _fig_prods = px.bar(
                                    _prod_rates, x="% with product", y="Product",
                                    orientation="h",
                                    title="Product Ownership Rates",
                                    color="% with product", color_continuous_scale="Purpor",
                                    text=_prod_rates["% with product"].map(lambda v: f"{v:.1f}%"),
                                )
                                _fig_prods.update_traces(textposition="outside")
                                _fig_prods.update_layout(
                                    height=340, margin=dict(l=20, r=80, t=40, b=20),
                                    coloraxis_showscale=False,
                                    xaxis=dict(ticksuffix="%", range=[0, 100]),
                                )
                                st.plotly_chart(_fig_prods, width='stretch')
                        with _ac7:
                            _nprod_cnt = (
                                _aud["n_products"].value_counts().sort_index()
                                .reset_index()
                            )
                            _nprod_cnt.columns = ["Products", "Customers"]
                            _fig_np = px.bar(
                                _nprod_cnt, x="Products", y="Customers",
                                title="# Products Held",
                                color="Customers", color_continuous_scale="YlOrRd",
                            )
                            _fig_np.add_scatter(
                                x=_nprod_cnt["Products"], y=_nprod_cnt["Customers"],
                                mode="lines+markers",
                                line=dict(color="#5a0000", width=2, shape="spline", smoothing=0.8),
                                marker=dict(size=6, color="#5a0000"),
                                showlegend=False, name="",
                            )
                            _fig_np.update_layout(
                                height=340, margin=dict(l=20, r=20, t=50, b=30),
                                coloraxis_showscale=False,
                            )
                            st.plotly_chart(_fig_np, width='stretch')

                        st.markdown("<br>", unsafe_allow_html=True)

                        # ── Row 4: Deposits vs IXI scatter (full width) ───────────
                        _sow_sc_df = _aud[["amount_deposit_spot_balance", "total_deposits_ixi", "sow"]].dropna() if "total_deposits_ixi" in _aud.columns else None
                        if _sow_sc_df is not None and len(_sow_sc_df) > 0:
                            _sow_sc_df = _sow_sc_df[
                                (_sow_sc_df["amount_deposit_spot_balance"] <= 200_000) &
                                (_sow_sc_df["total_deposits_ixi"] <= 1_500_000)
                            ]
                            _fig_sow = px.scatter(
                                _sow_sc_df,
                                x="amount_deposit_spot_balance",
                                y="total_deposits_ixi",
                                color="sow",
                                color_continuous_scale="RdYlGn",
                                range_color=[0, 1],
                                title="Deposits VS IXI",
                                labels={
                                    "amount_deposit_spot_balance": "Deposit ($)",
                                    "total_deposits_ixi": "IXI ($)",
                                    "sow": "SoW",
                                },
                                opacity=0.55,
                            )
                            _fig_sow.update_xaxes(range=[0, 200_000])
                            _fig_sow.update_yaxes(range=[0, 1_500_000])
                            # ── Percentile contour lines (add_shape layer='above' → always on top) ──
                            if len(_sow_sc_df) >= 20:
                                import plotly.graph_objects as go
                                import matplotlib
                                matplotlib.use('Agg')
                                import matplotlib.pyplot as _plts
                                from scipy.stats import gaussian_kde as _gkde_sow
                                from scipy.ndimage import gaussian_filter as _gfs
                                _xs = _sow_sc_df["amount_deposit_spot_balance"].values
                                _ys = _sow_sc_df["total_deposits_ixi"].values
                                try:
                                    _kdes = _gkde_sow(np.vstack([_xs, _ys]), bw_method=0.3)
                                    _xis = np.linspace(_xs.min(), _xs.max(), 200)
                                    _yis = np.linspace(_ys.min(), _ys.max(), 200)
                                    _XXs, _YYs = np.meshgrid(_xis, _yis)
                                    _ZZs = _kdes(np.vstack([_XXs.ravel(), _YYs.ravel()])).reshape(_XXs.shape)
                                    _ZZs = _gfs(_ZZs, sigma=3)
                                    _z_pts_s = _kdes(np.vstack([_xs, _ys]))
                                    for _plbls, _clrs in [
                                        (90, "#2166ac"),   # blue  — outermost (90% of pts inside)
                                        (50, "#1a9850"),   # green — 50% inside
                                        (25, "#f4a11d"),   # amber — 25% inside
                                        (10, "#c0392b"),   # red   — dense core
                                    ]:
                                        _lvls = float(np.percentile(_z_pts_s, 100 - _plbls))
                                        if _lvls <= 0:
                                            continue
                                        # Dummy trace for legend entry only
                                        _fig_sow.add_trace(go.Scatter(
                                            x=[None], y=[None], mode='lines',
                                            line=dict(color=_clrs, width=2.5),
                                            name=f'{_plbls}th pct',
                                            legendgroup=f'pct{_plbls}_s', showlegend=True,
                                        ))
                                        _mfigs, _maxs = _plts.subplots()
                                        _css = _maxs.contour(_XXs, _YYs, _ZZs, levels=[_lvls])
                                        _plts.close(_mfigs)
                                        for _segs in _css.get_paths():
                                            _vertss = _segs.vertices
                                            if len(_vertss) < 5:
                                                continue
                                            # Open contour at lower-left, extend to both axes
                                            _myis = int(np.argmin(_vertss[:, 1]))
                                            _mxis = int(np.argmin(_vertss[:, 0]))
                                            _ns = len(_vertss)
                                            _rols = np.roll(_vertss, -_myis, axis=0)
                                            _nmxs = (_mxis - _myis) % _ns
                                            _arc_as = _rols[:_nmxs + 1]
                                            _arc_bs = _rols[_nmxs:]
                                            if _arc_as.shape[0] < 3 or _arc_bs.shape[0] < 3:
                                                _opns = _rols
                                            elif np.mean(_arc_as[:, 0] + _arc_as[:, 1]) >= np.mean(_arc_bs[:, 0] + _arc_bs[:, 1]):
                                                _opns = _arc_as
                                            else:
                                                _opns = _arc_bs
                                            _ptss = np.vstack([
                                                [max(float(_opns[0, 0]), 0), 0],
                                                _opns,
                                                [0, max(float(_opns[-1, 1]), 0)],
                                            ])
                                            _svgs = "M " + " L ".join(f"{float(_px):.4f},{float(_py):.4f}" for _px, _py in _ptss)
                                            _fig_sow.add_shape(
                                                type='path', path=_svgs,
                                                xref='x', yref='y',
                                                line=dict(color=_clrs, width=2.5),
                                                fillcolor='rgba(0,0,0,0)',
                                                layer='above',
                                            )
                                except Exception:
                                    pass
                            _fig_sow.update_layout(height=380, margin=dict(l=20, r=20, t=50, b=30),
                                                   legend=dict(x=0.01, y=0.99, bgcolor='rgba(255,255,255,0.8)',
                                                               bordercolor='#aaa', borderwidth=1))
                            st.plotly_chart(_fig_sow, width='stretch')

            elif st.session_state.get("sim_run_triggered") and not _sim_segs_eff:
                st.warning("No segments to analyse — all included segments are also in the exclude list. Remove some exclusions.")


    # ══════════════════════════════════════════════════════════
    # DATA QUALITY TAB
    # ══════════════════════════════════════════════════════════
    with tab_audit:
        st.subheader("Data Quality & Recommendations")
        st.caption("Live health check of your dataset and current filter settings.")

        k1, k2, k3, k4 = st.columns(4)
        n_rows_raw   = len(df_raw)
        n_users_raw  = df_raw["alpha_key"].nunique() if "alpha_key" in df_raw.columns else 0
        n_segs       = df["nsegment"].nunique()
        contact_rate = df_raw["contact_flag"].mean() if "contact_flag" in df_raw.columns else 0
        k1.metric("Total records",       f"{n_rows_raw:,}")
        k2.metric("Unique customers",    f"{n_users_raw:,}")
        k3.metric("Unique segments",     f"{n_segs:,}")
        k4.metric("Contacted %",         f"{contact_rate:.1%}")

        st.divider()

        with st.expander("Data Quality", expanded=True):
            nan_bal  = df["balance_pct_change"].isna().mean()
            nan_acct = df["accounts_pct_change"].isna().mean()
            zero_bal = (df["start_balance"].fillna(0) < 1).mean()
            dup_pct  = df.duplicated(subset=["alpha_key", "communication", "nsegment"]).mean()
            dq1, dq2, dq3, dq4 = st.columns(4)
            dq1.metric("Balance data missing",    f"{nan_bal:.1%}")
            dq2.metric("Accounts data missing",   f"{nan_acct:.1%}")
            dq3.metric("Near-zero balances",       f"{zero_bal:.1%}")
            dq4.metric("Duplicate records",        f"{dup_pct:.1%}")
            date_min_v = df["start_date"].min()
            date_max_v = df["end_date"].max()
            st.caption(f"Date range: **{date_min_v.date() if pd.notna(date_min_v) else 'N/A'}** → **{date_max_v.date() if pd.notna(date_max_v) else 'N/A'}**")

        with st.expander("Methodology Notes", expanded=True):
            if bal_baseline_min is not None:
                st.info(f"**Low-balance filter active**: segments with avg starting balance < €{bal_baseline_min:,.0f} excluded from ranking.")

            st.markdown("**Control group coverage per communication:**")
            ctrl_cov = (
                df.groupby("communication")["contact_flag"]
                .apply(lambda x: (x == 0).mean())
                .reindex(COMM_ORDER).dropna().rename("Control %").to_frame()
            )
            ctrl_cov["Contact %"] = (1 - ctrl_cov["Control %"]).map("{:.1%}".format)
            ctrl_cov["Control %"] = ctrl_cov["Control %"].map("{:.1%}".format)
            st.dataframe(ctrl_cov, width='stretch')

            if "tbl" in dir() and not tbl.empty:
                n_cols_tbl = [c for c in tbl.columns if c.endswith("_n")]
                if n_cols_tbl:
                    low_n_count  = (tbl[n_cols_tbl] < 30).sum().sum()
                    total_cells  = tbl[n_cols_tbl].notna().sum().sum()
                    pct_low      = low_n_count / total_cells if total_cells else 0
                    msg = f"**Low-confidence cells (N < 30):** {low_n_count:,} of {total_cells:,} ({pct_low:.0%})."
                    (st.warning if pct_low > 0.3 else st.info)(msg)

            overlap     = df.groupby("alpha_key")["nsegment"].nunique()
            multi_pct   = (overlap > 1).mean()
            avg_segs_u  = overlap.mean()
            st.info(f"**Segment overlap:** {multi_pct:.1%} of customers in >1 segment (avg {avg_segs_u:.1f}/customer).")

        with st.expander("Top Segment Recommendations", expanded=True):
            if "tbl" not in dir() or tbl.empty:
                st.warning("No table data — adjust filters first.")
            else:
                def _top_tbl(metric_col: str, label: str) -> None:
                    if metric_col not in tbl.columns:
                        return
                    n_col = next((f"{c}_n" for c in ordered_comms if metric_col.startswith(c + "_")), None)
                    cols    = [metric_col] + ([n_col] if n_col and n_col in tbl.columns else [])
                    top     = tbl[cols].dropna(subset=[metric_col]).sort_values(metric_col, ascending=False).head(10).copy()
                    top.index = [f"{x}  —  {SEGMENT_LABELS.get(str(x), '')}" for x in top.index]
                    if n_col and n_col in top.columns:
                        top["Confidence"] = top[n_col].map(
                            lambda n: "🟢 High" if n >= 100 else ("🟡 Medium" if n >= 30 else "🔴 Low"))
                        top[n_col] = top[n_col].map("{:,.0f}".format)
                    top[metric_col] = top[metric_col].map("{:.2%}".format)
                    top.columns = ([label, "N", "Confidence"] if len(top.columns) == 3 else [label])
                    st.dataframe(top, width='stretch')

                _fc = ordered_comms[0]  if ordered_comms else None
                _lc = ordered_comms[-1] if ordered_comms else None
                rc1, rc2, rc3, rc4 = st.columns(4)
                with rc1:
                    st.markdown(f"**Top 10 — {_fc} Bal%**")
                    if _fc: _top_tbl(f"{_fc}_bal", "Bal% Δ")
                with rc2:
                    st.markdown(f"**Top 10 — {_lc} Bal%**")
                    if _lc: _top_tbl(f"{_lc}_bal", "Bal% Δ")
                with rc3:
                    st.markdown(f"**Top 10 — {_fc} Lift Bal**")
                    if _fc: _top_tbl(f"{_fc}_lift_bal", "Lift Bal")
                with rc4:
                    st.markdown(f"**Top 10 — {_fc} Lift Acct**")
                    if _fc: _top_tbl(f"{_fc}_lift_acct", "Lift Acct")


# ══════════════════════════════════════════════════════════
# CC BALANCE TRANSFER CAMPAIGN
# ══════════════════════════════════════════════════════════
elif active_campaign == "CC_BT":

    # ── Load raw CSV ──────────────────────────────────────────────────────────
    try:
        with open(_APP_DIR / "CC_BT_campaign.csv", "rb") as _f:
            _bt_raw_bytes = _f.read()
        _bt_raw = _load_cc_bt(_bt_raw_bytes)
    except FileNotFoundError:
        st.error("CC_BT_campaign.csv not found.")
        st.stop()

    _bt_required = ["alpha_key", "control_flag", "bt_amount", "bt_flag", "nsegments"]
    _bt_missing = [c for c in _bt_required if c not in _bt_raw.columns]
    if _bt_missing:
        st.error("CC_BT_campaign.csv is missing columns: " + ", ".join(_bt_missing))
        st.stop()

    _bt_all_seg_ids = _cc_bt_seg_ids(_bt_raw)

    # ── Sidebar ───────────────────────────────────────────────────────────────
    with st.sidebar:
        # ── Segment lookup ────────────────────────────────────────────────────
        _seg_lookup_widget(_bt_all_seg_ids, SEGMENT_LABELS, SEGMENT_DESCRIPTIONS)

        # ── Date range (shown only when bt_date column exists) ────────────────
        _bt_has_date = "bt_date" in _bt_raw.columns
        _bt_date_from = _bt_date_to = None
        if _bt_has_date:
            st.subheader("Date range")
            _bt_dates = pd.to_datetime(_bt_raw["bt_date"], errors="coerce").dropna()
            _bt_d_min = _bt_dates.min().date()
            _bt_d_max = _bt_dates.max().date()
            _bt_dcol1, _bt_dcol2 = st.columns(2)
            with _bt_dcol1:
                _bt_date_from = st.date_input("From", value=_bt_d_min, min_value=_bt_d_min,
                                              max_value=_bt_d_max, key="bt_date_from")
            with _bt_dcol2:
                _bt_date_to = st.date_input("To", value=_bt_d_max, min_value=_bt_d_min,
                                            max_value=_bt_d_max, key="bt_date_to")
            st.divider()

        with st.expander("Advanced settings", expanded=False):
            st.caption(
                "**Minimum N** — segments with fewer than **30** treated customers are hidden "
                "from the table. Too few observations make the mean unreliable and the CI very wide."
            )
            if _bt_has_date:
                st.caption(
                    "**Date range** — filters rows to customers whose balance transfer date "
                    "falls within the selected window."
                )

    # ── Explode & aggregate ───────────────────────────────────────────────────
    _bt_df  = preprocess_cc_bt(_bt_raw,
                                date_min=str(_bt_date_from) if _bt_date_from else None,
                                date_max=str(_bt_date_to) if _bt_date_to else None)
    _bt_tbl = agg_cc_bt(_bt_df, min_n=30)

    # ── Tabs ──────────────────────────────────────────────────────────────────
    _bt_tab_explorer, _bt_tab_data, _bt_tab_sim, _bt_tab_dq = st.tabs([
        "Segment Explorer",
        "Data",
        "Audience Simulator",
        "Data Quality",
    ])

    # ════ SEGMENT EXPLORER TAB ════════════════════════════════════════════════
    with _bt_tab_explorer:
        st.subheader("Segment catalogue")
        st.caption("Browse all segments by BT conversion performance. Unique Customers = treated customers in this campaign.")
        if not SEGMENT_LABELS:
            st.info("No segment descriptions found. Ensure segment_descriptions.csv is present.")
        else:
            _bt_exp_base = pd.DataFrame({
                "Segment ID": list(SEGMENT_LABELS.keys()),
                "Group":      list(SEGMENT_LABELS.values()),
                "Description": [SEGMENT_DESCRIPTIONS.get(k, "") for k in SEGMENT_LABELS],
            })
            _bt_user_counts = (
                _bt_df[_bt_df["control_flag"] == 0]
                .groupby("nsegment")["alpha_key"]
                .nunique()
                .rename("Unique Customers")
                .reset_index()
                .rename(columns={"nsegment": "Segment ID"})
            )
            _bt_exp_base = _bt_exp_base.merge(_bt_user_counts, on="Segment ID", how="left")
            _bt_exp_base["Unique Customers"] = _bt_exp_base["Unique Customers"].fillna(0).astype(int)
            _bt_exp_base = _bt_exp_base.sort_values("Segment ID").reset_index(drop=True)
            _bt_exp_styler = (
                _bt_exp_base.set_index("Segment ID").rename_axis(None)
                .style
                .apply(lambda s: s.map(_n_color), subset=["Unique Customers"], axis=0)
                .format(na_rep="")
            )
            _bt_exp_h = max(300, min(700, 60 + len(_bt_exp_base) * 26))
            components.html(
                _styled_html_table(_bt_exp_styler, SEGMENT_LABELS, SEGMENT_DESCRIPTIONS, height=_bt_exp_h),
                height=_bt_exp_h, scrolling=True,
            )
            st.divider()
            _grp_sz_bt = (
                _bt_exp_base.groupby("Group")
                .agg(Segments=("Segment ID", "count"), Customers=("Unique Customers", "sum"))
                .reset_index()
                .sort_values("Customers", ascending=False)
            )
            _fig_gsz_bt = px.bar(
                _grp_sz_bt, x="Group", y="Segments",
                color="Segments", color_continuous_scale="Blues",
                text="Segments", title="Segments per group",
            )
            _fig_gsz_bt.update_xaxes(tickangle=-30)
            _fig_gsz_bt.update_traces(textposition="outside")
            _fig_gsz_bt.update_layout(
                coloraxis_showscale=False, height=360,
                yaxis_range=[0, int(_grp_sz_bt["Segments"].max()) * 1.20],
            )
            st.plotly_chart(_fig_gsz_bt, width='stretch')
            _fig_gus_bt = px.bar(
                _grp_sz_bt, x="Group", y="Customers",
                color="Customers", color_continuous_scale="Teal",
                text=_grp_sz_bt["Customers"].map(lambda v: f"{v:,}"),
                title="Unique treated customers per group",
            )
            _fig_gus_bt.update_xaxes(tickangle=-30)
            _fig_gus_bt.update_traces(textposition="outside")
            _fig_gus_bt.update_layout(
                coloraxis_showscale=False, height=360,
                yaxis_range=[0, int(_grp_sz_bt["Customers"].max()) * 1.20],
            )
            st.plotly_chart(_fig_gus_bt, width='stretch')

    # ════ DATA TAB ════════════════════════════════════════════════════════════
    with _bt_tab_data:
        if _bt_tbl.empty:
            st.warning("No segments meet the minimum N threshold. Lower the filter in the sidebar.")
        else:
            # ── Metric / display options ──────────────────────────────────────
            _bt_mc1, _bt_mc2, _ = st.columns([2, 2, 3])
            _bt_metric = _bt_mc1.radio(
                "Metric", ["BT Conversion Lift", "BT Amount Lift"],
                horizontal=True, key="bt_metric_radio",
            )
            _bt_show_n = _bt_mc2.checkbox("Show N", value=False, key="bt_show_n")

            if _bt_metric == "BT Conversion Lift":
                _val_col, _ci_col  = "conv_lift", "conv_lift_ci"
                _t_col,   _c_col   = "conv_treated", "conv_control"
                _val_label = "Conv Lift"
                _t_label   = "Conv % (treated)"
                _c_label   = "Conv % (control)"
                _fmt_fn    = lambda v: f"{v:.2%}" if pd.notna(v) else ""
                _lo, _hi   = -0.05, 0.15
            else:
                _val_col, _ci_col  = "amt_lift", "amt_lift_ci"
                _t_col,   _c_col   = "amt_treated", "amt_control"
                _val_label = "Amt Lift ($)"
                _t_label   = "Avg BT $ (treated)"
                _c_label   = "Avg BT $ (control)"
                _fmt_fn    = lambda v: f"${v:,.0f}" if pd.notna(v) else ""
                _lo, _hi   = -500, 2000

            # Build display table — only lift + CI (+ N when requested)
            _disp_cols = [_val_col, _ci_col]
            if _bt_show_n:
                _disp_cols = ["n_treated", "n_control"] + _disp_cols
            _bt_disp = _bt_tbl[[c for c in _disp_cols if c in _bt_tbl.columns]].copy()
            _bt_disp.index = [str(x) for x in _bt_disp.index]
            _col_rename = {
                "n_treated": "N (treated)", "n_control": "N (control)",
                _val_col: _val_label, _ci_col: "±CI",
            }
            _bt_disp = _bt_disp.rename(columns=_col_rename)

            # Colour the lift column
            def _bt_colour(series):
                return series.map(
                    lambda v: _rdylgn(v, lo=_lo, hi=_hi)
                    if isinstance(v, (int, float)) and not pd.isna(v) else ""
                )

            _fmt_map = {_val_label: _fmt_fn, "±CI": _fmt_fn}
            _styled_bt = _bt_disp.style.format(_fmt_map, na_rep="")
            if "N (treated)" in _bt_disp.columns:
                _styled_bt = _styled_bt.format(
                    {"N (treated)": "{:,.0f}", "N (control)": "{:,.0f}"}, na_rep=""
                )
            if _val_label in _bt_disp.columns:
                _styled_bt = _styled_bt.apply(_bt_colour, subset=[_val_label], axis=0)

            _bt_h = max(300, min(800, 60 + len(_bt_disp) * 26))
            components.html(
                _styled_html_table(_styled_bt, height=_bt_h),
                height=_bt_h, scrolling=True,
            )

            # ── Downloads (right below table, before chart) ───────────────────
            _bt_metric_key = "conv" if _bt_metric == "BT Conversion Lift" else "amt"
            _, _bt_btn_col, _ = st.columns([2, 3, 2])
            _bt_btn_col.download_button(
                label="Download Excel",
                data=build_cc_bt_excel(_bt_tbl, _bt_metric_key),
                file_name="cc_bt_segment_lift.xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                use_container_width=True,
                key="bt_dl_xlsx",
            )

            # ── Recommended Audiences ────────────────────────────────────────
            st.divider()
            st.markdown("### Recommended Audiences")
            st.caption("Optimal segment groupings for the BT offer, ranked by lift. "
                       "Use these as ready-made targeting lists — no manual segment selection needed.")
            _bt_ra_c1, _bt_ra_c2, _bt_ra_c3 = st.columns([3, 3, 1])
            with _bt_ra_c1:
                _bt_ra_metric = st.radio("Optimise for", ["Conv Lift", "Amt Lift"],
                                         horizontal=True, key="bt_ra_metric")
            with _bt_ra_c2:
                _bt_ra_all_opts = sorted(_bt_tbl.index.astype(str).tolist())
                _bt_ra_and_segs = st.multiselect(
                    "AND — mandatory segments", options=_bt_ra_all_opts,
                    default=[s for s in st.session_state.get("bt_ra_and_segs", []) if s in _bt_ra_all_opts],
                    key="bt_ra_and_segs")
            with _bt_ra_c3:
                _bt_ra_min_aud = st.number_input("Min audience", min_value=0, value=500, step=100, key="bt_ra_min_aud")
            _bt_ra_val_col = "conv_lift" if _bt_ra_metric == "Conv Lift" else "amt_lift"
            _bt_ra_ci_col  = "conv_lift_ci" if _bt_ra_metric == "Conv Lift" else "amt_lift_ci"
            _bt_ra_label   = "Conv Lift" if _bt_ra_metric == "Conv Lift" else "Amt Lift ($)"
            _bt_ra_lo      = -0.05 if _bt_ra_metric == "Conv Lift" else -500
            _bt_ra_hi      = 0.15  if _bt_ra_metric == "Conv Lift" else 2000
            _bt_ra_fmt     = ((lambda v: f"{v:.2%}" if pd.notna(v) else "")
                              if _bt_ra_metric == "Conv Lift"
                              else (lambda v: f"${v:,.0f}" if pd.notna(v) else ""))
            _bt_ra_placeholder = st.empty()
            _bt_ra_shown_segs: List[str] = []
            _bt_ra_bot_segs:   List[str] = []
            _bt_and_idx_ra:    List[str] = []
            if _bt_ra_val_col in _bt_tbl.columns:
                _bt_ra_placeholder.info("⧗ Building recommended audience…")
                _bt_naive = _bt_tbl[_bt_ra_val_col].dropna()
                _bt_naive.index = _bt_naive.index.astype(str)
                _bt_sorted_ra = _bt_naive.sort_values(ascending=False)
                _bt_treated_all = _bt_df[_bt_df["control_flag"] == 0]
                _bt_and_str = set(str(s) for s in (_bt_ra_and_segs or []))
                if _bt_and_str:
                    _bt_csg = (_bt_treated_all[_bt_treated_all["nsegment"].astype(str).isin(_bt_and_str)]
                               .groupby("alpha_key")["nsegment"].apply(lambda x: set(x.astype(str))))
                    _bt_valid = _bt_csg[_bt_csg.apply(lambda s: _bt_and_str.issubset(s))].index
                    _bt_treated_filt = _bt_treated_all[_bt_treated_all["alpha_key"].isin(_bt_valid)]
                else:
                    _bt_treated_filt = _bt_treated_all
                _bt_sc = {str(s): set(_bt_treated_filt[_bt_treated_filt["nsegment"].astype(str) == str(s)]["alpha_key"])
                           for s in _bt_sorted_ra.index}
                _bt_and_idx_ra = [s for s in (_bt_ra_and_segs or []) if s in _bt_sorted_ra.index]
                _bt_non_and   = _bt_sorted_ra[~_bt_sorted_ra.index.isin(_bt_and_idx_ra)]
                _bt_bot_n     = max(3, min(15, max(1, len(_bt_non_and)) // 4))
                _bt_ra_bot_ser = _bt_non_and.sort_values().head(_bt_bot_n)
                _bt_not_strs  = set(str(s) for s in _bt_ra_bot_ser.index)
                _bt_not_custs_ra: set = set()
                for _bns in _bt_ra_bot_ser.index:
                    _bt_not_custs_ra |= _bt_sc.get(str(_bns), set())
                _bt_run_ra: set = ((set().union(*[_bt_sc[str(s)] for s in _bt_and_idx_ra]) - _bt_not_custs_ra)
                                   if _bt_and_idx_ra else set())
                _bt_sel_ra: List[str] = []
                for _bts in _bt_non_and.index:
                    if str(_bts) in _bt_not_strs:
                        continue
                    if len(_bt_run_ra) >= _bt_ra_min_aud:
                        break
                    _bt_sel_ra.append(_bts)
                    _bt_run_ra |= (_bt_sc[str(_bts)] - _bt_not_custs_ra)
                _bt_ra_top_idx = list(dict.fromkeys(_bt_and_idx_ra + _bt_sel_ra))
                _bt_ra_top = _bt_sorted_ra.reindex([s for s in _bt_ra_top_idx if s in _bt_sorted_ra.index]).dropna()
                _bt_final_custs = _bt_run_ra
                _bt_ra_users = len(_bt_final_custs)
                _bt_tc = _bt_treated_filt[_bt_treated_filt["alpha_key"].isin(_bt_final_custs)].drop_duplicates("alpha_key")
                _bt_cc = _bt_df[(_bt_df["control_flag"] == 1) &
                                (_bt_df["nsegment"].astype(str).isin([str(s) for s in _bt_ra_top_idx]))].drop_duplicates("alpha_key")
                if _bt_ra_metric == "Conv Lift":
                    _bt_tv  = _bt_tc["bt_flag"].mean() if not _bt_tc.empty else np.nan
                    _bt_cv_ = _bt_cc["bt_flag"].mean() if not _bt_cc.empty else np.nan
                else:
                    _bt_tv  = _bt_tc[_bt_tc["bt_flag"] == 1]["bt_amount"].mean() if not _bt_tc.empty else np.nan
                    _bt_cv_ = _bt_cc[_bt_cc["bt_flag"] == 1]["bt_amount"].mean() if not _bt_cc.empty else np.nan
                _bt_w_lift = float(_bt_tv - _bt_cv_) if (pd.notna(_bt_tv) and pd.notna(_bt_cv_)) else np.nan
                _bt_rm1, _bt_rm2 = st.columns(2)
                _bt_rm1.metric("Recommended audience", f"{_bt_ra_users:,} customers")
                _bt_rm2.metric(f"Expected {_bt_ra_label}", _bt_ra_fmt(_bt_w_lift),
                               help="Treated audience mean minus matched-control mean — each customer counted once.")
                def _bt_make_ra_table(series):
                    _ci_v2 = (_bt_tbl[_bt_ra_ci_col].rename(index=str).reindex(series.index.astype(str)).values
                              if _bt_ra_ci_col in _bt_tbl.columns else [np.nan] * len(series))
                    _n_v2  = [int(_bt_tbl.at[s, "n_treated"]) if "n_treated" in _bt_tbl.columns and s in _bt_tbl.index else 0
                              for s in series.index]
                    _d2 = pd.DataFrame({
                        "Description": [SEGMENT_DESCRIPTIONS.get(str(s), SEGMENT_LABELS.get(str(s), "")) for s in series.index],
                        _bt_ra_label: series.values,
                        "±CI": _ci_v2,
                        "N": _n_v2,
                    }, index=series.index)
                    _d2.index.name = None
                    return (_d2.style
                            .format({_bt_ra_label: _bt_ra_fmt, "±CI": _bt_ra_fmt, "N": "{:,.0f}"}, na_rep="")
                            .apply(lambda s: s.map(lambda v: _rdylgn(v, lo=_bt_ra_lo, hi=_bt_ra_hi)),
                                   subset=[_bt_ra_label], axis=0)
                            .apply(lambda s: s.map(_n_color), subset=["N"], axis=0))
                _bt_ra_sty     = _bt_make_ra_table(_bt_ra_top)
                _bt_ra_bot_sty = _bt_make_ra_table(_bt_ra_bot_ser)
                _bt_ra_shown_segs = [str(s) for s in _bt_ra_top.index]
                _bt_ra_bot_segs   = [str(s) for s in _bt_ra_bot_ser.index]
                _bt_comb_h = max(300, min(900, 80 + (len(_bt_ra_top) + len(_bt_ra_bot_ser)) * 30))
                components.html(
                    _two_tables_html([
                        (f"Top {len(_bt_ra_top)} segments (highest lift)", _bt_ra_sty, [(_bt_ra_label, "±CI")]),
                        (f"Bottom {len(_bt_ra_bot_ser)} segments (lowest lift — suppressed)", _bt_ra_bot_sty, [(_bt_ra_label, "±CI")]),
                    ], _bt_comb_h),
                    height=_bt_comb_h, scrolling=True,
                )
                _bt_ra_placeholder.empty()
            else:
                _bt_ra_placeholder.empty()
                st.info(f"No {_bt_ra_label} data available.")
            if _bt_ra_shown_segs:
                _, _bt_ra_btn_col, _ = st.columns([1, 2, 1])
                if _bt_ra_btn_col.button("Send to Audience Simulator", key="bt_ra_send_sim", use_container_width=True):
                    st.session_state["bt_sim_segs"]          = _bt_ra_shown_segs
                    st.session_state["bt_sim_segs_and"]      = _bt_and_idx_ra
                    st.session_state["bt_sim_segs_excl"]     = _bt_ra_bot_segs
                    st.session_state["bt_sim_run_triggered"] = True
                    st.success(f"Sent {len(_bt_ra_shown_segs)} segments to Audience Simulator "
                               f"(+ {len(_bt_ra_bot_segs)} excluded).")
        # ── Audience Demographics (all customers in audience_profile.csv) ─────
        if _aud_df is not None and len(_aud_df) > 0:
            st.divider()
            st.subheader("Audience Demographics — All Customers")
            st.caption(f"{len(_aud_df):,} customers in audience_profile.csv")
            import plotly.graph_objects as go
            from scipy.stats import gaussian_kde as _gkde_dt

            def _kde_dt(series, nbins, line_color):
                _v = series.dropna().values
                if len(_v) < 5:
                    return None
                _, _edges = np.histogram(_v, bins=nbins)
                _bw = _edges[1] - _edges[0]
                _fn = _gkde_dt(_v)
                _xs = np.linspace(_v.min(), _v.max(), 300)
                _ys = _fn(_xs) * len(_v) * _bw
                return go.Scatter(x=_xs, y=_ys, mode="lines",
                                  line=dict(color=line_color, width=2.5),
                                  showlegend=False, name="")

            _d = _aud_df.copy()
            _d1, _d2 = st.columns([3, 2])
            with _d1:
                _age_bins2   = [18, 25, 35, 45, 55, 65, 75, 120]
                _age_labels2 = ["18-24", "25-34", "35-44", "45-54", "55-64", "65-74", "75+"]
                _d["_age_group2"] = pd.cut(_d["age"], bins=_age_bins2, labels=_age_labels2, right=False)
                _age_dist2 = (_d["_age_group2"].value_counts().reindex(_age_labels2).fillna(0).reset_index())
                _age_dist2.columns = ["Age group", "Customers"]
                _fig_age2 = px.bar(_age_dist2, x="Age group", y="Customers",
                                   title="Age Distribution", color="Customers",
                                   color_continuous_scale="Blues")
                _fig_age2.add_scatter(x=_age_dist2["Age group"], y=_age_dist2["Customers"],
                                      mode="lines+markers",
                                      line=dict(color="#103060", width=2, shape="spline", smoothing=1.0),
                                      marker=dict(size=6, color="#103060"), showlegend=False, name="")
                _age_mean2 = float(_d["age"].dropna().mean())
                _age_mean_bin2 = str(pd.cut([_age_mean2], bins=_age_bins2, labels=_age_labels2, right=False)[0])
                _fig_age2.add_shape(type="line", xref="x", yref="paper",
                                    x0=_age_mean_bin2, x1=_age_mean_bin2, y0=0, y1=1,
                                    line=dict(color="red", width=2, dash="dash"))
                _fig_age2.add_annotation(x=_age_mean_bin2, yref="paper", y=1.05,
                                         text=f"Mean: {_age_mean2:.1f}y", showarrow=False,
                                         xanchor="left", font=dict(color="red", size=11))
                _fig_age2.update_layout(height=300, margin=dict(l=20, r=20, t=50, b=30),
                                        coloraxis_showscale=False, showlegend=False)
                st.plotly_chart(_fig_age2, width='stretch')
            with _d2:
                _gender_counts2 = _d["gender"].fillna("Missing").value_counts().reset_index()
                _gender_counts2.columns = ["Gender", "Count"]
                _g_cmap2 = {"Male": "#89C4E1", "Female": "#FFB7C5", "Missing": "#CCCCCC"}
                _fig_gender2 = go.Figure(go.Pie(
                    labels=_gender_counts2["Gender"], values=_gender_counts2["Count"],
                    hole=0.4,
                    marker=dict(colors=[_g_cmap2.get(g, "#DDDDDD") for g in _gender_counts2["Gender"]]),
                    textinfo="percent+label"))
                _fig_gender2.update_layout(title=dict(text="Gender"), height=300,
                                           margin=dict(l=20, r=20, t=50, b=30))
                st.plotly_chart(_fig_gender2, width='stretch')
            st.markdown("<br>", unsafe_allow_html=True)
            _d3, _d4, _d5 = st.columns(3)
            with _d3:
                _fig_ten2 = px.histogram(_d, x="tenure_years", nbins=20, title="Tenure Distribution",
                                         labels={"tenure_years": "Tenure (years)"},
                                         color_discrete_sequence=["#4C9BE8"])
                _k2 = _kde_dt(_d["tenure_years"], 20, "#1a3a6b")
                if _k2:
                    _fig_ten2.add_trace(_k2)
                _ten_med2 = float(_d["tenure_years"].dropna().median())
                _fig_ten2.add_vline(x=_ten_med2, line=dict(color="red", width=2, dash="dash"),
                                    annotation_text=f"Median: {_ten_med2:.1f}y",
                                    annotation_position="top right",
                                    annotation_font=dict(color="red", size=11))
                _fig_ten2.update_layout(height=300, margin=dict(l=20, r=20, t=50, b=30))
                st.plotly_chart(_fig_ten2, width='stretch')
            with _d4:
                _dep_p99 = float(_d["amount_deposit_spot_balance"].dropna().quantile(0.99))
                _fig_dep2 = px.histogram(_d[_d["amount_deposit_spot_balance"] <= _dep_p99],
                                         x="amount_deposit_spot_balance", nbins=25, title="Deposit Balance",
                                         labels={"amount_deposit_spot_balance": "Balance ($)"},
                                         color_discrete_sequence=["#F4A261"])
                _k3 = _kde_dt(_d["amount_deposit_spot_balance"], 25, "#7a3100")
                if _k3:
                    _fig_dep2.add_trace(_k3)
                _dep_med2 = float(_d["amount_deposit_spot_balance"].dropna().median())
                _fig_dep2.add_vline(x=_dep_med2, line=dict(color="red", width=2, dash="dash"),
                                    annotation_text=f"Median: ${_dep_med2:,.0f}",
                                    annotation_position="top right",
                                    annotation_font=dict(color="red", size=11))
                _fig_dep2.update_xaxes(range=[0, _dep_p99])
                _fig_dep2.update_layout(height=300, margin=dict(l=20, r=20, t=50, b=30))
                st.plotly_chart(_fig_dep2, width='stretch')
            with _d5:
                _state_cnt2 = _d["state"].value_counts().head(10).reset_index()
                _state_cnt2.columns = ["State", "Customers"]
                _fig_state2 = px.bar(_state_cnt2, x="Customers", y="State", orientation="h",
                                     title="Top 10 States", color="Customers", color_continuous_scale="Teal")
                _fig_state2.update_layout(height=300, margin=dict(l=20, r=20, t=50, b=30),
                                          coloraxis_showscale=False,
                                          yaxis=dict(categoryorder="total ascending"))
                st.plotly_chart(_fig_state2, width='stretch')
            st.markdown("<br>", unsafe_allow_html=True)
            _d6, _d7 = st.columns(2)
            with _d6:
                _prod_cols_dt = [c for c in PRODUCT_COLS if c in _d.columns]
                if _prod_cols_dt:
                    _prod_rates2 = ((_d[_prod_cols_dt].mean() * 100).sort_values(ascending=True).reset_index())
                    _prod_rates2.columns = ["Product", "% with product"]
                    _fig_prods2 = px.bar(_prod_rates2, x="% with product", y="Product", orientation="h",
                                         title="Product Ownership Rates", color="% with product",
                                         color_continuous_scale="Purpor",
                                         text=_prod_rates2["% with product"].map(lambda v: f"{v:.1f}%"))
                    _fig_prods2.update_traces(textposition="outside")
                    _fig_prods2.update_layout(height=340, margin=dict(l=20, r=80, t=40, b=20),
                                              coloraxis_showscale=False,
                                              xaxis=dict(ticksuffix="%", range=[0, 100]))
                    st.plotly_chart(_fig_prods2, width='stretch')
            with _d7:
                _nprod_cnt2 = _d["n_products"].value_counts().sort_index().reset_index()
                _nprod_cnt2.columns = ["Products", "Customers"]
                _fig_np2 = px.bar(_nprod_cnt2, x="Products", y="Customers", title="# Products Held",
                                   color="Customers", color_continuous_scale="YlOrRd")
                _fig_np2.add_scatter(x=_nprod_cnt2["Products"], y=_nprod_cnt2["Customers"],
                                      mode="lines+markers",
                                      line=dict(color="#5a0000", width=2, shape="spline", smoothing=0.8),
                                      marker=dict(size=6, color="#5a0000"), showlegend=False, name="")
                _nprod_mean2 = float(_d["n_products"].dropna().mean())
                _fig_np2.add_vline(x=_nprod_mean2, line=dict(color="red", width=2, dash="dash"),
                                   annotation_text=f"Mean: {_nprod_mean2:.1f}",
                                   annotation_position="top right",
                                   annotation_font=dict(color="red", size=11))
                _fig_np2.update_layout(height=340, margin=dict(l=20, r=20, t=50, b=30),
                                        coloraxis_showscale=False)
                st.plotly_chart(_fig_np2, width='stretch')

    # ════ AUDIENCE SIMULATOR TAB ════════════════════════════════════════════════
    with _bt_tab_sim:
        st.subheader("Audience Performance Simulator")
        st.caption(
            "Select any set of segments and we'll estimate the **expected BT conversion rate "
            "and incremental revenue** if you targeted only those customers. "
            "Each customer is counted once — the lift shown is what you'd actually observe on that cohort."
        )
        if _bt_tbl.empty:
            st.warning("No table data — adjust the minimum-N filter in the sidebar first.")
        else:
            _bt_sim_all = sorted(_bt_tbl.index.astype(str).tolist())
            _bt_sim_or = st.multiselect(
                "OR — segments to include (union)", options=_bt_sim_all,
                default=[s for s in st.session_state.get("bt_sim_segs", []) if s in _bt_sim_all],
                key="bt_sim_segs",
                help="A customer qualifies if they belong to ANY of these segments.",
            )
            _bt_sim_and = st.multiselect(
                "AND — must also be in (intersection)", options=_bt_sim_all,
                default=[s for s in st.session_state.get("bt_sim_segs_and", []) if s in _bt_sim_all],
                key="bt_sim_segs_and",
                help="Only keeps customers in both OR and AND lists.",
            )
            _bt_sim_excl = st.multiselect(
                "NOT — segments to exclude", options=_bt_sim_all,
                default=[s for s in st.session_state.get("bt_sim_segs_excl", []) if s in _bt_sim_all],
                key="bt_sim_segs_excl",
                help="Removes these segments from the final audience.",
            )
            _bt_sim_eff = [s for s in (_bt_sim_or or [])
                           if (not _bt_sim_and or s in _bt_sim_and) and s not in (_bt_sim_excl or [])]
            if st.button("▶ Run simulation", key="bt_sim_run_btn", type="primary"):
                st.session_state["bt_sim_run_triggered"] = True
                st.session_state["_bt_sim_segs_snap"] = list(_bt_sim_eff)
                st.session_state["_bt_sim_excl_snap"] = list(_bt_sim_excl or [])
            _bt_compute_segs = st.session_state.get("_bt_sim_segs_snap", [])
            _bt_compute_excl = st.session_state.get("_bt_sim_excl_snap", [])
            if st.session_state.get("bt_sim_run_triggered") and _bt_compute_segs:
                _bt_not_ids_sim: set = set(
                    _bt_df[(_bt_df["control_flag"] == 0) &
                           (_bt_df["nsegment"].astype(str).isin(_bt_compute_excl))]["alpha_key"]
                ) if _bt_compute_excl else set()
                _bt_users_base = _bt_df[
                    (_bt_df["control_flag"] == 0) & (_bt_df["nsegment"].astype(str).isin(_bt_compute_segs))
                ]["alpha_key"]
                if _bt_not_ids_sim:
                    _bt_users_base = _bt_users_base[~_bt_users_base.isin(_bt_not_ids_sim)]
                _bt_total_users_sim = _bt_users_base.nunique()
                _bt_tot_c1, _bt_tot_c2 = st.columns([2, 1])
                _excl_note_sim = f" ({len(_bt_compute_excl)} excluded)" if _bt_compute_excl else ""
                _bt_tot_c1.metric(f"Total unique customers (all selected segments){_excl_note_sim}",
                                  f"{_bt_total_users_sim:,}")
                _bt_sim_metric = _bt_tot_c2.radio("Metric", ["Conv Lift", "Amt Lift"],
                                                  horizontal=True, key="bt_sim_metric")
                _bt_sim_treated = (_bt_df[(_bt_df["control_flag"] == 0) &
                                          (_bt_df["nsegment"].astype(str).isin(_bt_compute_segs))]
                                   .drop_duplicates("alpha_key"))
                if _bt_not_ids_sim:
                    _bt_sim_treated = _bt_sim_treated[~_bt_sim_treated["alpha_key"].isin(_bt_not_ids_sim)]
                _bt_sim_control = (_bt_df[(_bt_df["control_flag"] == 1) &
                                          (_bt_df["nsegment"].astype(str).isin(_bt_compute_segs))]
                                   .drop_duplicates("alpha_key"))
                _bt_rate_t = _bt_sim_treated["bt_flag"].mean() if not _bt_sim_treated.empty else np.nan
                _bt_rate_c = _bt_sim_control["bt_flag"].mean() if not _bt_sim_control.empty else np.nan
                _bt_lift_s = float(_bt_rate_t - _bt_rate_c) if (pd.notna(_bt_rate_t) and pd.notna(_bt_rate_c)) else np.nan
                _bt_conv_t = _bt_sim_treated[_bt_sim_treated["bt_flag"] == 1]
                _bt_avg_amt_s  = _bt_conv_t["bt_amount"].mean() if not _bt_conv_t.empty else 0.0
                _bt_incr_conv_s = int(len(_bt_sim_treated) * max(0.0, _bt_lift_s)) if pd.notna(_bt_lift_s) else 0
                _bt_incr_amt_s  = _bt_incr_conv_s * _bt_avg_amt_s
                st.markdown("#### Expected lift across selected segments")
                _sk1, _sk2, _sk3, _sk4 = st.columns(4)
                _sk1.metric("Audience size",         f"{len(_bt_sim_treated):,}")
                _sk2.metric("BT rate (treated)",     f"{_bt_rate_t:.1%}" if pd.notna(_bt_rate_t) else "—",
                            delta=f"{_bt_lift_s:+.1%} vs control" if pd.notna(_bt_lift_s) else "")
                _sk3.metric("Avg BT amount",         f"${_bt_avg_amt_s:,.0f}")
                _sk4.metric("Projected incr. conv.", f"{_bt_incr_conv_s:,}")
                st.markdown("#### Projected absolute impact")
                _sp1, _sp2, _sp3 = st.columns(3)
                _sp1.metric("Incremental BT amount", f"${_bt_incr_amt_s:,.0f}" if _bt_incr_amt_s else "—",
                            help="Projected converters × avg BT amount.")
                _sp2.metric("BT rate (control)",     f"{_bt_rate_c:.1%}" if pd.notna(_bt_rate_c) else "—")
                _sp3.metric("BT rate lift",          f"{_bt_lift_s:.1%}" if pd.notna(_bt_lift_s) else "—")
                st.divider()
                _bt_h_col_sim, _ = st.columns([3, 1])
                _bt_h_col_sim.markdown("#### Per-segment breakdown")
                _bt_sim_metric_key = "conv" if _bt_sim_metric == "Conv Lift" else "amt"
                _bt_intl_cols_sim = [c for c in
                    [f"{_bt_sim_metric_key}_lift", f"{_bt_sim_metric_key}_lift_ci", "n_treated"]
                    if c in _bt_tbl.columns]
                _bt_all_detail = list(dict.fromkeys(
                    list(_bt_sim_or or []) + list(_bt_sim_and or []) + list(_bt_sim_excl or [])))
                _bt_detail = _bt_tbl.loc[
                    _bt_tbl.index.astype(str).isin(_bt_all_detail),
                    [c for c in _bt_intl_cols_sim if c in _bt_tbl.columns],
                ].copy()
                _bt_detail.index.name = None
                _bt_lift_lbl_sim = "Conv Lift" if _bt_sim_metric == "Conv Lift" else "Amt Lift ($)"
                _bt_detail = _bt_detail.rename(columns={
                    f"{_bt_sim_metric_key}_lift": _bt_lift_lbl_sim,
                    f"{_bt_sim_metric_key}_lift_ci": "±CI",
                    "n_treated": "N (treated)",
                })
                _bt_sim_excl_set = set(str(s) for s in (_bt_sim_excl or []))
                _bt_sim_and_set  = set(str(s) for s in (_bt_sim_and  or []))
                def _bt_get_role(seg):
                    s = str(seg)
                    if s in _bt_sim_excl_set: return "NOT"
                    if s in _bt_sim_and_set:  return "AND"
                    return "OR"
                _bt_detail.insert(0, "Role", [_bt_get_role(s) for s in _bt_detail.index])
                _bt_detail = (_bt_detail.assign(_rs=_bt_detail["Role"].map({"AND": 0, "OR": 1, "NOT": 2}))
                              .sort_values("_rs").drop(columns="_rs"))
                _bt_fmt_sim = ((lambda v: f"{v:.2%}" if pd.notna(v) else "")
                               if _bt_sim_metric == "Conv Lift"
                               else (lambda v: f"${v:,.0f}" if pd.notna(v) else ""))
                _bt_sty_sim = _bt_detail.style.format(
                    {c: _bt_fmt_sim for c in [_bt_lift_lbl_sim, "±CI"] if c in _bt_detail.columns}
                    | ({"N (treated)": "{:,.0f}"} if "N (treated)" in _bt_detail.columns else {}),
                    na_rep="",
                )
                if _bt_lift_lbl_sim in _bt_detail.columns:
                    _lo_sim = -0.05 if _bt_sim_metric == "Conv Lift" else -500
                    _hi_sim = 0.15  if _bt_sim_metric == "Conv Lift" else 2000
                    _bt_sty_sim = _bt_sty_sim.apply(
                        lambda s: s.map(lambda v: _rdylgn(v, lo=_lo_sim, hi=_hi_sim)),
                        subset=[_bt_lift_lbl_sim], axis=0)
                def _bt_role_color_sim(v):
                    if v == "NOT": return "background-color: #ffdddd; color: #880000"
                    if v == "AND": return "background-color: #e8d5ff; color: #5500aa"
                    if v == "OR":  return "background-color: #d0e8ff; color: #003399"
                    return ""
                _bt_sty_sim = _bt_sty_sim.apply(lambda s: s.map(_bt_role_color_sim), subset=["Role"], axis=0)
                _bt_sim_det_h = max(300, min(600, 60 + len(_bt_detail) * 30))
                components.html(
                    _styled_html_table(_bt_sty_sim, SEGMENT_LABELS, SEGMENT_DESCRIPTIONS, height=_bt_sim_det_h),
                    height=_bt_sim_det_h, scrolling=True,
                )
                st.markdown("---")
                st.markdown("#### Audience Profile")
                if _aud_df is None:
                    st.info("No audience profile data found. Please provide an audience_profile.csv file.")
                else:
                    _bt_aud_keys2 = _bt_df[
                        (_bt_df["control_flag"] == 0) & (_bt_df["nsegment"].astype(str).isin(_bt_compute_segs))
                    ]["alpha_key"]
                    if _bt_compute_excl:
                        _bt_aud_excl2 = set(_bt_df[
                            (_bt_df["control_flag"] == 0) & (_bt_df["nsegment"].astype(str).isin(_bt_compute_excl))
                        ]["alpha_key"])
                        _bt_aud_keys2 = _bt_aud_keys2[~_bt_aud_keys2.isin(_bt_aud_excl2)]
                    _bt_aud2 = _aud_df[_aud_df["alpha_key"].isin(_bt_aud_keys2.unique())].copy()
                    if len(_bt_aud2) == 0:
                        st.info("No demographic records for the selected segments.")
                    else:
                        _bpk1, _bpk2, _bpk3, _bpk4 = st.columns(4)
                        _bpk1.metric("Customers with profile data", f"{len(_bt_aud2):,}",
                                     help="Customers from audience_profile.csv matched to this audience. May be a subset of the full campaign audience.")
                        _bpk2.metric("Median age",    f"{_bt_aud2['age'].median():.0f} yrs" if "age" in _bt_aud2.columns else "—")

                        _bpk3.metric("Median tenure", f"{_bt_aud2['tenure_years'].median():.1f} yrs" if "tenure_years" in _bt_aud2.columns else "—")
                        _bpk4.metric("Median SoW",    f"{_bt_aud2['sow'].median()*100:.0f}%" if "sow" in _bt_aud2.columns else "—")
                        import plotly.graph_objects as _go_bt2
                        _bac1, _bac2 = st.columns([3, 2])
                        with _bac1:
                            if "age" in _bt_aud2.columns:
                                _bt_ab = [18,25,35,45,55,65,75,120]
                                _bt_al = ["18-24","25-34","35-44","45-54","55-64","65-74","75+"]
                                _bt_aud2["_ag2"] = pd.cut(_bt_aud2["age"], bins=_bt_ab, labels=_bt_al, right=False)
                                _bt_agd = _bt_aud2["_ag2"].value_counts().reindex(_bt_al).fillna(0).reset_index()
                                _bt_agd.columns = ["Age group","Customers"]
                                _fig_ba2 = px.bar(_bt_agd, x="Age group", y="Customers",
                                                  title="Age Distribution", color="Customers",
                                                  color_continuous_scale="Blues")
                                _fig_ba2.update_layout(height=300, coloraxis_showscale=False)
                                st.plotly_chart(_fig_ba2, width='stretch')
                        with _bac2:
                            if "gender" in _bt_aud2.columns:
                                _bt_gc2 = _bt_aud2["gender"].fillna("Missing").value_counts().reset_index()
                                _bt_gc2.columns = ["Gender","Count"]
                                _fig_bg2 = _go_bt2.Figure(_go_bt2.Pie(
                                    labels=_bt_gc2["Gender"], values=_bt_gc2["Count"], hole=0.4,
                                    marker=dict(colors=["#89C4E1","#FFB7C5","#CCCCCC"]),
                                    textinfo="percent+label"))
                                _fig_bg2.update_layout(title="Gender", height=300)
                                st.plotly_chart(_fig_bg2, width='stretch')
                        _bac3, _bac4 = st.columns(2)
                        with _bac3:
                            if "tenure_years" in _bt_aud2.columns:
                                _fig_bten2 = px.histogram(_bt_aud2, x="tenure_years", nbins=20,
                                                           title="Tenure Distribution",
                                                           color_discrete_sequence=["#4C9BE8"])
                                _fig_bten2.update_layout(height=300)
                                st.plotly_chart(_fig_bten2, width='stretch')
                        with _bac4:
                            if "amount_deposit_spot_balance" in _bt_aud2.columns:
                                _fig_bdep2 = px.histogram(_bt_aud2, x="amount_deposit_spot_balance",
                                                            nbins=25, title="Deposit Balance",
                                                            color_discrete_sequence=["#F4A261"])
                                _fig_bdep2.update_layout(height=300)
                                st.plotly_chart(_fig_bdep2, width='stretch')
            elif st.session_state.get("bt_sim_run_triggered") and not _bt_sim_eff:
                st.warning("No segments to analyse — all included segments are also in the exclude list.")

    # ════ DATA QUALITY TAB ════════════════════════════════════════════════════
    with _bt_tab_dq:
        st.subheader("Data Quality & Recommendations")
        st.caption("Live health check of the CC Balance Transfer dataset and current filter settings.")

        _dq1, _dq2, _dq3, _dq4 = st.columns(4)
        _dq1.metric("Total records",    f"{len(_bt_raw):,}")
        _dq2.metric("Unique customers", f"{_bt_raw['alpha_key'].nunique():,}")
        _dq3.metric("Unique segments",  f"{_bt_df['nsegment'].nunique():,}")
        _dq4.metric("Treated %",        f"{(_bt_raw['control_flag'] == 0).mean():.1%}")

        st.divider()

        with st.expander("Data Quality", expanded=True):
            _nan_bt  = _bt_raw["bt_flag"].isna().mean()
            _nan_amt = _bt_raw["bt_amount"].isna().mean()
            _dup_bt  = _bt_raw.duplicated(subset=["alpha_key"]).mean()
            _bt_conv = (_bt_raw["bt_flag"] == 1).mean()
            _ddq1, _ddq2, _ddq3, _ddq4 = st.columns(4)
            _ddq1.metric("BT flag missing",     f"{_nan_bt:.1%}")
            _ddq2.metric("BT amount missing",   f"{_nan_amt:.1%}")
            _ddq3.metric("BT conversion rate",  f"{_bt_conv:.1%}")
            _ddq4.metric("Duplicate customers", f"{_dup_bt:.1%}")

        with st.expander("Methodology Notes", expanded=True):
            _ctrl_split_bt  = (_bt_raw["control_flag"] == 1).mean()
            _treat_split_bt = (_bt_raw["control_flag"] == 0).mean()
            st.markdown(
                f"**Treated / Control split:** {_treat_split_bt:.1%} treated — {_ctrl_split_bt:.1%} control"
            )
            _bt_overlap = (_bt_df.groupby("alpha_key")["nsegment"].nunique() > 1).mean()
            _bt_avg_segs = _bt_df.groupby("alpha_key")["nsegment"].nunique().mean()
            st.info(
                f"**Segment overlap:** {_bt_overlap:.1%} of customers in >1 segment "
                f"(avg {_bt_avg_segs:.1f}/customer)."
            )
            _bt_seg_dist = _bt_df.groupby("nsegment")["alpha_key"].nunique().describe()
            st.caption(
                f"Customers per segment — min: **{int(_bt_seg_dist['min']):,}** | "
                f"median: **{int(_bt_seg_dist['50%']):,}** | "
                f"max: **{int(_bt_seg_dist['max']):,}**"
            )

        with st.expander("Top Segment Recommendations", expanded=True):
            if _bt_tbl.empty:
                st.warning("No table data — adjust filters first.")
            else:
                _top_conv_bt = (
                    _bt_tbl[["conv_lift", "n_treated"]].dropna(subset=["conv_lift"])
                    .sort_values("conv_lift", ascending=False)
                    .head(10)
                    .copy()
                )
                _top_conv_bt["Confidence"] = _top_conv_bt["n_treated"].map(
                    lambda n: "🟢 High" if n >= 100 else ("🟡 Medium" if n >= 30 else "🔴 Low")
                )
                _top_conv_bt.index = [
                    f"{x}  —  {SEGMENT_LABELS.get(str(x), '')}" for x in _top_conv_bt.index
                ]
                _top_conv_bt = _top_conv_bt.rename(
                    columns={"conv_lift": "Conv Lift", "n_treated": "N (treated)"}
                )
                _top_conv_bt["Conv Lift"] = _top_conv_bt["Conv Lift"].map("{:.2%}".format)
                _top_conv_bt["N (treated)"] = _top_conv_bt["N (treated)"].map("{:,.0f}".format)
                _rc1_bt, _rc2_bt = st.columns(2)
                with _rc1_bt:
                    st.markdown("**Top 10 segments — BT Conversion Lift**")
                    st.dataframe(_top_conv_bt, width='stretch')
                with _rc2_bt:
                    if "amt_lift" in _bt_tbl.columns:
                        _top_amt_bt = (
                            _bt_tbl[["amt_lift", "n_treated"]].dropna(subset=["amt_lift"])
                            .sort_values("amt_lift", ascending=False)
                            .head(10)
                            .copy()
                        )
                        _top_amt_bt.index = [
                            f"{x}  —  {SEGMENT_LABELS.get(str(x), '')}" for x in _top_amt_bt.index
                        ]
                        _top_amt_bt = _top_amt_bt.rename(
                            columns={"amt_lift": "Amt Lift ($)", "n_treated": "N (treated)"}
                        )
                        _top_amt_bt["Amt Lift ($)"] = _top_amt_bt["Amt Lift ($)"].map("${:,.0f}".format)
                        _top_amt_bt["N (treated)"] = _top_amt_bt["N (treated)"].map("{:,.0f}".format)
                        st.markdown("**Top 10 segments — BT Amount Lift**")
                        st.dataframe(_top_amt_bt, width='stretch')
